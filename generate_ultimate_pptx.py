import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------
# Brand Theme Color Palette (Hexa Force Dark Mode)
# ----------------------------------------------------
C_BG = RGBColor(15, 23, 42)        # Slate 900 (Dark Slate Background)
C_CARD = RGBColor(30, 41, 59)      # Slate 800 (Card Panel Background)
C_PRIMARY = RGBColor(37, 99, 235)  # Electric Blue (Primary Headings)
C_TEXT = RGBColor(241, 245, 249)   # Cool White (Body Text)
C_MUTED = RGBColor(148, 163, 184)  # Muted Gray (Subtitles/Sidebars)
C_RED = RGBColor(239, 68, 68)      # Crimson Red (Attacker/Vulnerable Callouts)
C_GREEN = RGBColor(16, 185, 129)   # Emerald Green (Defender/Hardened Callouts)
C_WHITE = RGBColor(255, 255, 255)  # Pure White

# ----------------------------------------------------
# Slide Layout Helper Functions
# ----------------------------------------------------
def apply_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG

def add_header(slide, title, category="HEXA FORCE CONTAINER SECURITY SERIES"):
    # Category Tag
    cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_PRIMARY
    
    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(11.83), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE

def draw_card(slide, left, top, width, height, bg_color=C_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_card_text(slide, left, top, width, height, title, body_bullets, title_color=C_WHITE, font_size_body=13):
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Card Title
    if title:
        p_title = tf.paragraphs[0]
        p_title.text = title.upper()
        p_title.font.name = "Arial"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.space_after = Pt(10)
        
    for bullet in body_bullets:
        p = tf.add_paragraph() if tf.paragraphs[-1].text else tf.paragraphs[0]
        p.text = bullet
        p.font.name = "Arial"
        p.font.size = Pt(font_size_body)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(6)
        
        # Simple indentation formatting if bullet starts with indent character or bullet character
        if bullet.strip().startswith("•") or bullet.strip().startswith("-"):
            p.level = 0
        else:
            p.level = 0

def add_title_slide(slide, main_title, subtitle, category="HEXA FORCE DEVSECOPS RESEARCH"):
    apply_dark_bg(slide)
    
    # Structural Accent Band
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(0.12), Inches(3.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_PRIMARY
    accent.line.fill.background()
    
    # Category
    cat_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.0), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.0), Inches(1.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = main_title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(42)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.0), Inches(1.0))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = C_MUTED

# ----------------------------------------------------
# Main Data Generators
# ----------------------------------------------------
def build_deck_1(prs):
    # DECK 1 slides
    slides_d1 = [
        # Slide 1: Title
        {"type": "title", "title": "HEXA FORCE CONTAINER LAB", "subtitle": "4-Stage DevSecOps Attack & Mitigation Laboratory\nCryptographically Signed Container Architectures, OIDC Trust, and DevSecOps Containment Rules", "cat": "Hexa Force Lab Delivery"},
        
        # Slide 2: Agenda
        {"type": "content2", "title": "Briefing & Laboratory Agenda",
         "c1_title": "Section 1: Vulnerability & Exploitation", "c1_bullets": ["• Core Architecture: Shared Linux kernel virtualization and standard container security assumptions.", "• The Exploit: Detailed execution of the CVE-2016-5195 Copy-on-Write local privilege escalation.", "• 4-Stage Attack Chain: Moving from low-privileged shell to process capability abuse, API exposure, and persistent host takeover."],
         "c2_title": "Section 2: Defense-in-Depth & Pipeline", "c2_bullets": ["• Runtime Mitigation: Implementing strict Linux capability profiles, non-root system users, and read-only container rootfs.", "• Supply Chain Provenance: Constructing a secure CI/CD build pipeline using Cosign and GitHub Actions OpenID Connect (OIDC).", "• Integrity Validation: Keyless container signing verified against public transparency ledgers (Rekor)."]},
        
        # Slide 3: 4 Pillars
        {"type": "content2", "title": "The 4 Security Pillars of Containers",
         "c1_title": "Pillar 1 & 2: Kernel & Process Isolation", "c1_bullets": ["• Host Kernel Sharing: The core lightweight efficiency of containers is their chief weakness: direct exposure to local host kernel exploits.", "• Linux Capabilities: Processes are limited by partitioned system call permissions. Excess privileges (like CAP_SYS_ADMIN) break container namespaces."],
         "c2_title": "Pillar 3 & 4: API & Volume Boundaries", "c2_bullets": ["• Daemon API Security: Mounting '/var/run/docker.sock' directly inside a container allows the container root user to command the host Docker engine.", "• Filesystem & Volume Mounts: Poor volume mount scopes allow path traversals, filesystem write access, and persistent injection into host runtimes."]},
        
        # Slide 4: Stage 1
        {"type": "content2", "title": "Stage 1: Host Kernel Isolation (Dirty COW)",
         "c1_title": "Exploit Vector: CVE-2016-5195", "c1_bullets": ["• Initial shell secured as unprivileged user 'victim' inside the Docker container.", "• Exploit compiled: 'dirtyc0w.c' uses memory mapping and multi-threaded race conditions.", "• Memory target: Overwriting a read-only file ('/var/secret/target.txt') owned by container root."],
         "c2_title": "Technical Outcome", "c2_bullets": ["• The exploit utilizes the 'madvise(MADV_DONTNEED)' system call to discard mapped page frames.", "• Concurrently writes payload through the '/proc/self/mem' interface.", "• Result: Bypasses file permissions, overwrites target, and escalates process privileges to container root."]},
        
        # Slide 5: Stage 2
        {"type": "content2", "title": "Stage 2: Process & Capabilities Isolation",
         "c1_title": "Vulnerability Analysis", "c1_bullets": ["• Linux capabilities partition superuser powers into distinct, assignable units.", "• If a container is misconfigured with 'CAP_SYS_ADMIN', its processes gain massive control.", "• This includes namespace management and raw storage device manipulation."],
         "c2_title": "Exploitation Flow", "c2_bullets": ["• The newly elevated container root user queries active capabilities.", "• Finding CAP_SYS_ADMIN active, the attacker bypasses process limits.", "• This privilege forms the architectural beachhead needed to compromise external host interfaces."]},
        
        # Slide 6: Stage 3
        {"type": "content2", "title": "Stage 3: Daemon & API Security (docker.sock)",
         "c1_title": "Threat Vector: Exposed Unix Sockets", "c1_bullets": ["• Continuous integration and monitoring agents often mount the host's Docker socket '/var/run/docker.sock' inside the container.", "• This provides the container process with a direct control channel to the host's physical Docker daemon."],
         "c2_title": "Command Execution", "c2_bullets": ["• Since the attacker is now root inside the container, they gain read/write permissions to the mounted socket.", "• The attacker issues direct Docker API curl commands to the daemon.", "• They command the host to pull external images and instantiate new containers."]},
        
        # Slide 7: Stage 4
        {"type": "content2", "title": "Stage 4: Volume & Filesystem Security",
         "c1_title": "Physical Breakout Vector", "c1_bullets": ["• Using the exposed Docker socket, the attacker spawns a new privileged container on the host.", "• Configuration: The host's physical root filesystem ('/') is mounted directly to '/mnt/host' inside the container.", "• The attacker now has write permissions to all physical host files."],
         "c2_title": "Host Takeover & Persistence", "c2_bullets": ["• The attacker navigates into '/mnt/host/etc/cron.d' and injects a custom cron job.", "• The cron job triggers a root-privileged reverse shell back to the attack system.", "• Physical host takeover is complete, establishing a persistent footprint on the node."]},
        
        # Slide 8: Supply Chain Image Signing
        {"type": "content2", "title": "Cryptographic Image Signing (Cosign & OIDC)",
         "c1_title": "The Keyless Trust Paradigm", "c1_bullets": ["• Traditional container signing depends on managing sensitive private keys in CI/CD environments. If leaked, the entire supply chain falls.", "• Hexa Force implements modern keyless signing via Sigstore/Cosign, eliminating key storage."],
         "c2_title": "Cryptographic Architecture", "c2_bullets": ["• The runner requests a short-lived X.509 certificate from the Fulcio Certificate Authority.", "• Authentication is federated using an OpenID Connect (OIDC) identity token issued by GitHub Actions.", "• Signatures are logged in the immutable Rekor public transparency ledger."]},
        
        # Slide 9: Pipeline Orchestration
        {"type": "content2", "title": "Automated CI/CD Laboratory Pipeline",
         "c1_title": "Automated DevSecOps Workflow", "c1_bullets": ["• Source code push triggers the GitHub Actions workflow runner.", "• Docker image compiled and tagged dynamically with lowercase parameter expansion.", "• Image pushed securely to the GitHub Container Registry (GHCR)."],
         "c2_title": "Automated Verification Gates", "c2_bullets": ["• Immediate execution of Cosign keyless signing using the runner's OIDC identity.", "• High-performance validation gates verify signatures during cluster deployment.", "• Result: Completely blocks unsigned or unverified images from executing in production."]},
        
        # Slide 10: Best Practices Hardening
        {"type": "content2", "title": "DevSecOps Hardening & Defense-in-Depth",
         "c1_title": "Runtime Defenses", "c1_bullets": ["• Drop excessive Linux capabilities: drop CAP_SYS_ADMIN and run with '--cap-drop=ALL' by default.", "• Set read-only filesystems: run containers with '--read-only' to block exploit downloads and compilations.", "• Non-root execution: enforce 'USER victim' or non-root UID policies."],
         "c2_title": "Host & Infrastructure Defenses", "c2_bullets": ["• Host Kernel Patching: Strictly enforce regular patching schedules for CVE-2016-5195 and other kernel exploits.", "• Socket Restrictions: Completely block mounting '/var/run/docker.sock' in application containers.", "• LSM Sandboxing: Deploy AppArmor or SELinux profiles."]},
        
        # Slide 11: Key Takeaways
        {"type": "content2", "title": "Conclusion: Hexa Force Lab Debriefing",
         "c1_title": "Summary of Success", "c1_bullets": ["• Exploit chain fully mapped: demonstrated how minor misconfigurations interact with kernel vulnerabilities to allow host escape.", "• Supply Chain Secured: implemented modern keyless signing to guarantee deployment provenance."],
         "c2_title": "Strategic Mission", "c2_bullets": ["• Security as an Enabler: integrating cryptographic signing directly into CI/CD ensures automation without slowing delivery.", "• Continuous Hardening: container virtualization is not a security boundary; defense must be applied continuously across the pipeline."]}
    ]
    
    # ----------------------------------------------------
    # DECK 1 Speaker Notes Content Mapped Slide-by-Slide
    # ----------------------------------------------------
    d1_notes_text = [
        # Slide 1
        "=== SPOKEN TALK TRACK ===\n"
        "\"Good morning, evaluators and distinguished guests. Welcome to the technical delivery of the Hexa Force "
        "Container Security Lab. Today, we are presenting our end-to-end container security demonstration, modeling a "
        "high-fidelity 4-stage attack chain and corresponding defense-in-depth hardening strategies. Over the next few "
        "minutes, we will walk you through how a local privilege escalation vulnerability can cascade into complete host "
        "takeover, and how modern cryptographic verification using keyless Sigstore/Cosign signing with GitHub Actions OIDC "
        "robustly secures the software delivery pipeline. Let us begin by reviewing our core objectives.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Establish the Hexa Force brand immediately with an authoritative and professional delivery.\n"
        "• Set the scope of the presentation: it is a balanced demonstration of an advanced attack chain and enterprise-grade cryptographic defense.\n"
        "• Emphasize key modern pipeline terms: Sigstore, Cosign, OIDC federation, and GHCR registry integrations.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What motivated the choice of this specific attack chain?\n"
        "Answer: We chose this specific progression because it realistically models how simple runtime misconfigurations and unpatched host-level vulnerabilities interact in enterprise microservices. It demonstrates that container security cannot rely on runtime isolation alone, but must encompass host hardening, image provenance, and CI/CD integrity.",
        
        # Slide 2
        "=== SPOKEN TALK TRACK ===\n"
        "\"Our presentation today is structured into four main areas. First, we establish the threat landscape and the "
        "4 Security Pillars of Containers. Next, we walk step-by-step through our 4-stage attack simulation, showcasing how "
        "we break kernel isolation, processes, APIs, and volume mounts. Third, we present our cryptographic supply chain "
        "integrity pipeline using Cosign and OpenID Connect. Finally, we conclude with our enterprise-grade hardening "
        "recommendations and DevSecOps best practices to secure any modern production cluster.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Guide the audience through the agenda to show structural clarity and excellent preparation.\n"
        "• Emphasize that the lab moves logically from threat identification (simulation) to threat neutralization (hardening).\n"
        "• Make it clear that the defense-in-depth framework applies to all levels: runtime, pipeline, and host kernel.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the most critical phase of the agenda?\n"
        "Answer: While the exploit phase demonstrates the weakness, the supply chain integrity section is the most critical. Securing the container runtime is only half the battle; we must also guarantee code provenance through modern cryptographic signing to prevent unauthorized container execution.",
        
        # Slide 3
        "=== SPOKEN TALK TRACK ===\n"
        "\"To secure any containerized environment, we must evaluate the four foundational security pillars: (1) Host Kernel "
        "Isolation, (2) Process & Capabilities Isolation, (3) Daemon & API Security, and (4) Volume & Filesystem Security. "
        "Standard container systems share the host kernel, which means a kernel vulnerability exposes the host. Processes "
        "inside a container must be strictly limited via Linux capabilities to prevent privilege abuse. The Docker API daemon "
        "socket must never be exposed to untrusted environments, and physical volume mounts must be strictly governed to "
        "prevent directory traversal and host escapes. Let us look at Stage 1, where kernel isolation is broken.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Introduce the 4 Pillars as the industry-standard frame of reference for evaluating container security.\n"
        "• Explain the shared kernel model: containers are lightweight because they share the host kernel, which is also their primary single point of failure.\n"
        "• Emphasize that security must be applied across all 4 pillars simultaneously; securing only one layer leaves the system vulnerable.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why is kernel sharing the ultimate double-edged sword for containers?\n"
        "Answer: Soapbox container systems gain lightweight speed and resource efficiency from sharing the kernel, but inherit a unified single point of failure. If the host kernel has a vulnerability, container barriers are fundamentally bypassed, collapsing the first and most critical isolation pillar.",
        
        # Slide 4
        "=== SPOKEN TALK TRACK ===\n"
        "\"At Stage 1, we execute our first exploit inside an unprivileged container. Running as the unprivileged user "
        "'victim', we execute our compiled binary, 'dirtyc0w'. This exploit targets the Dirty COW kernel vulnerability, "
        "which is a race condition in the kernel's virtual memory subsystem's handling of Copy-on-Write. We target a "
        "read-only secret file, '/var/secret/target.txt'. By using madvise with the MADV_DONTNEED flag, we race the kernel's "
        "write mechanisms, tricking the kernel into writing root-owned data into a read-only memory segment. This instantly "
        "elevates our privilege from an unprivileged user to container root.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Focus heavily on the virtual memory mechanics: Copy-on-Write, madvise system call, and the O_RDONLY file bypass.\n"
        "• Clearly declare CVE-2016-5195 to demonstrate deep technical mastery and precise vulnerability tracking.\n"
        "• Emphasize that this first exploit is executed as a non-root, unprivileged user, simulating a compromised microservice shell.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why does a local privilege escalation inside a container matter if it is isolated?\n"
        "Answer: Local privilege escalation is the beachhead. While namespaces still limit what a container root user can see, obtaining root privileges is the prerequisite for abusing other exposed interfaces, such as Linux capabilities (like CAP_SYS_ADMIN), raw network sockets, or system files.",
        
        # Slide 5
        "=== SPOKEN TALK TRACK ===\n"
        "\"In Stage 2, we exploit process-level privileges. Having achieved container root in Stage 1, we now check "
        "the container's capabilities. If a container is run with elevated Linux capabilities, such as CAP_SYS_ADMIN, "
        "the boundaries of standard container isolation begin to disintegrate. Linux capabilities partition traditional "
        "superuser root powers into distinct, assignable units. By abusing CAP_SYS_ADMIN, we gain the authority to modify "
        "kernel namespaces and mount external devices, laying the physical groundwork for our escape to the host filesystem.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Explain the concept of Linux capabilities: splitting traditional all-or-nothing root powers into fine-grained permissions.\n"
        "• Emphasize CAP_SYS_ADMIN as the root of all process isolation failures inside containers.\n"
        "• Discuss the principle of least privilege: containers must drop all unnecessary capabilities by default.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What are some alternatives to granting CAP_SYS_ADMIN to containers?\n"
        "Answer: Administrators should perform strict capability profiling and grant only the narrowest possible capabilities (like CAP_NET_BIND_SERVICE for low-port binding) instead of CAP_SYS_ADMIN. If system-level tasks are required, a dedicated daemon or sidecar pattern should be evaluated.",
        
        # Slide 6
        "=== SPOKEN TALK TRACK ===\n"
        "\"Stage 3 demonstrates a critical lateral movement vector: the exposed Docker socket. In many development "
        "environments, administrators mount the host's Docker socket—'/var/run/docker.sock'—directly into the container "
        "to enable nested docker tasks. Since our Stage 1 exploit gave us root privileges inside the container, we have "
        "full read/write permission to this socket. Because the Docker socket is the direct command API to the host's "
        "Docker engine, we can now send API requests directly to the host daemon, giving us full command over the host's "
        "container lifecycle.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Highlight the absolute danger of the Docker socket (/var/run/docker.sock) exposure.\n"
        "• Describe the Docker socket as the entry point to the host's physical container engine.\n"
        "• Explain how controlling the socket allows the containerized root user to send administrative API commands directly to the host daemon.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the principal security risk of mounting the Docker socket inside a container?\n"
        "Answer: Mounting the Docker socket is equivalent to granting root access to the physical host. Anyone with access to the socket can send commands to the host Docker daemon to pull arbitrary images, create new privileged containers, and mount the host's physical hard drives, rendering container isolation completely obsolete.",
        
        # Slide 7
        "=== SPOKEN TALK TRACK ===\n"
        "\"In Stage 4, we execute the final breakout. Using our access to the Docker socket, we issue an API command to "
        "launch a new, highly privileged container on the host. In its configuration, we mount the host's physical root "
        "filesystem to '/mnt/host' inside the container. Because this container runs with full host namespace access, "
        "we traverse directly into '/mnt/host/etc/cron.d' and write a malicious cron job. When the host's cron daemon "
        "executes, it runs our script and opens a root-privileged reverse shell back to our attack machine. We have achieved "
        "complete host takeover.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Detail the mechanical steps of host escape: privileged container instantiation $\rightarrow$ root filesystem mount $\rightarrow$ cron job injection.\n"
        "• Emphasize the threat of persistence: writing to `/etc/cron.d` guarantees that the host daemon will re-execute the reverse shell periodically.\n"
        "• Establish that this completes a full, high-severity exploit chain from a low-privileged container user to total physical host control.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How can host filesystem mounting be blocked at the enterprise level?\n"
        "Answer: Enterprises should implement Kubernetes Pod Security Admissions (PSA) or open-source policies (like OPA Gatekeeper or Kyverno) to strictly enforce that containers cannot run in privileged mode, share the host PID/IPC namespaces, or mount host paths.",
        
        # Slide 8
        "=== SPOKEN TALK TRACK ===\n"
        "\"Beyond runtime hardening, Hexa Force secures the software supply chain using state-of-the-art keyless "
        "container signing via Cosign. Traditional signing relies on storing highly sensitive cryptographic private keys. "
        "If these keys are leaked, the supply chain is compromised. Hexa Force avoids this completely. We leverage OpenID "
        "Connect to exchange GitHub Action build identities for temporary, 10-minute certificates. Cosign signs the built "
        "image, recording the signature in the immutable public Rekor transparency log, guaranteeing container origin and integrity.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Contrast traditional key management with modern keyless container signing.\n"
        "• Define the role of Sigstore: Fulcio issues short-lived certificates, Rekor provides the immutable public ledger.\n"
        "• Explain the cryptographic trust path: the signature is bound to a specific OIDC builder identity, verifying code provenance.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is Fulcio's role in keyless signing?\n"
        "Answer: Fulcio is the Sigstore Certificate Authority. It issues short-lived (10-minute) X.509 certificates. The certificate is bound to an OIDC token issued by a trusted identity provider (like GitHub Actions), proving that the build pipeline itself signed the container.",
        
        # Slide 9
        "=== SPOKEN TALK TRACK ===\n"
        "\"We operationalize this supply chain defense by automating the entire pipeline inside GitHub Actions. Upon a push "
        "to the main branch, our workflow triggers. It builds the Docker image, dynamically converts the image tag to "
        "lowercase using bash expansion to ensure compatibility with container registry specifications, and pushes it to "
        "the GitHub Container Registry. Immediately, Cosign uses the OIDC identity to sign the image. This seamless, "
        "automated pipeline guarantees that only verified and signed code is allowed to run in our production environment.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Highlight the automation aspects of the GHA workflow: building, lowercasing, pushing to GHCR, and signing.\n"
        "• Explain the lowercase string translation `${GITHUB_REPOSITORY,,}` as a crucial Docker OCI registry compliance step.\n"
        "• Position the CI/CD pipeline as the core policy enforcement point of your security program.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How does signature verification work inside a Kubernetes cluster?\n"
        "Answer: A validation controller (like Cosign Policy Controller or Kyverno) runs as an Admission Controller in the cluster. When a new pod request is made, the controller queries the Rekor transparency log to verify the image's cryptographic signature against the allowed OIDC builder identity, rejecting any unsigned or unverified images.",
        
        # Slide 10
        "=== SPOKEN TALK TRACK ===\n"
        "\"Let us examine how we completely remediate this threat. Hexa Force implements a comprehensive Defense-in-Depth "
        "architecture. First, we restrict container capabilities by explicitly dropping CAP_SYS_ADMIN. Second, we run the "
        "container as a non-root user and mark the root filesystem as read-only. Third, we block all exposure of the Docker "
        "socket. Finally, we apply kernel security patches to eliminate Dirty COW at the host level. Even if an attacker "
        "gains an initial shell, these combined controls prevent them from ever gaining root or escaping the container.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Emphasize the four layers of defense: capability dropping, non-root runtimes, read-only rootfs, and kernel patching.\n"
        "• Discuss the concept of Defense-in-Depth: even if one layer fails, subsequent layers prevent the exploit from succeeding.\n"
        "• Position these controls as standard, low-overhead enterprise recommendations.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the security value of a read-only container root filesystem?\n"
        "Answer: A read-only root filesystem prevents an attacker from writing files, installing tools, downloading exploit scripts, or compiling C binaries (like `dirtyc0w.c`) inside the container. It dramatically reduces the available attack surface and prevents runtime environmental tampering.",
        
        # Slide 11
        "=== SPOKEN TALK TRACK ===\n"
        "\"In conclusion, the Hexa Force Container Security Lab successfully demonstrates the extreme vulnerability of "
        "default container runtimes, while delivering an enterprise-ready blueprint for comprehensive defense. By pairing "
        "strict runtime hardening policies with cryptographic supply chain verification, we ensure complete resilience "
        "against host compromises and supply chain injection. Thank you for your time. The Hexa Force team is now open "
        "to your technical questions.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Deliver a clean, professional closing statement summarizing the core achievements.\n"
        "• Reiterate the value of combining runtime isolation with cryptographic verification.\n"
        "• Invite questions with absolute confidence, ready to leverage your detailed Q&As.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the single most important lesson from this research?\n"
        "Answer: Security is a multi-layered process. We cannot rely on container virtualization as a bulletproof sandbox. We must harden the host kernel, enforce the principle of least privilege inside the container, and use cryptographic pipelines to guarantee that only trusted code ever enters the execution environment."
    ]
    
    # Render Deck 1 Slide Designs
    blank_layout = prs.slide_layouts[6] # Blank
    
    for i, slide_data in enumerate(slides_d1):
        slide = prs.slides.add_slide(blank_layout)
        apply_dark_bg(slide)
        
        if slide_data["type"] == "title":
            add_title_slide(slide, slide_data["title"], slide_data["subtitle"], slide_data["cat"])
        else:
            add_header(slide, slide_data["title"])
            
            # Draw Two Columns
            col_w = Inches(5.6)
            col_h = Inches(4.8)
            top_y = Inches(1.8)
            
            # Left Card
            draw_card(slide, Inches(0.75), top_y, col_w, col_h)
            add_card_text(slide, Inches(0.75), top_y, col_w, col_h, slide_data["c1_title"], slide_data["c1_bullets"], title_color=C_PRIMARY)
            
            # Right Card (If hardening, outline with Emerald Green; if attack, outline with Slate)
            border_c = C_GREEN if "Hardening" in slide_data["title"] or "Remediation" in slide_data["title"] or "Signing" in slide_data["title"] or "Takeaways" in slide_data["title"] else None
            draw_card(slide, Inches(6.8), top_y, col_w, col_h, border_color=border_c)
            add_card_text(slide, Inches(6.8), top_y, col_w, col_h, slide_data["c2_title"], slide_data["c2_bullets"], title_color=C_GREEN if border_c else C_PRIMARY)
            
        # Inject Speaker Notes
        notes_slide = slide.notes_slide
        tf_notes = notes_slide.notes_text_frame
        tf_notes.text = d1_notes_text[i]
        
    print("[+] Generated Deck 1 Slides programmatically.")

def build_deck_2(prs):
    # DECK 2 slides
    slides_d2 = [
        # Slide 1: Title
        {"type": "title", "title": "DIRTY COW EXPLOIT DEEP DIVE", "subtitle": "Vulnerability Analysis, Virtual Memory Mechanics, and Container Escapes (CVE-2016-5195)\nPresented by Hexa Force Security Research", "cat": "Dirty COW Briefing"},
        
        # Slide 2: Agenda
        {"type": "content2", "title": "Briefing & Laboratory Agenda",
         "c1_title": "Foundational Operating System Concepts", "c1_bullets": ["• Virtual Memory: address translations, page tables, and physical RAM management.", "• Copy-On-Write (COW): operating system performance optimization routines.", "• The Vulnerability: analyzing the non-atomic race condition in core memory page fault structures."],
         "c2_title": "Exploitation & Mitigation", "c2_bullets": ["• The Exploit: multi-threaded race logic via '/proc/self/mem' and 'madvise'.", "• Container Risk: analyzing why shared host kernels collapse Docker isolation barriers.", "• Remediation: details of the Linux kernel 'FOLL_COW' patch and enterprise best practices."]},
        
        # Slide 3: Executive Summary
        {"type": "content2", "title": "Executive Summary: CVE-2016-5195",
         "c1_title": "Vulnerability Profile", "c1_bullets": ["• Core ID: CVE-2016-5195, widely known as the 'Dirty COW' kernel exploit.", "• Discovered: October 2016 by Phil Oester, affecting Linux kernel versions from 2.6.22 (2007) to 2016.", "• CVSS Severity: Rated High (CVSS 7.8) due to direct local privilege escalation."],
         "c2_title": "The Core Threat", "c2_bullets": ["• Bypasses absolute file access limits: allows unprivileged processes to write to read-only files.", "• Threat vectors: local users can overwrite system configurations (such as /etc/passwd) to gain root.", "• Impact: works natively within unprivileged Docker containers, compromising hosts."]},
         
        # Slide 4: Virtual Memory
        {"type": "content2", "title": "Linux Memory Management Foundations",
         "c1_title": "Virtual & Physical Memory Spaces", "c1_bullets": ["• Virtual Memory addresses isolate running processes from physical RAM hardware blocks.", "• Address Translation: The CPU's Memory Management Unit (MMU) uses page tables to translate pages.", "• Memory is mapped in fixed-size blocks (typically 4KB page frames)."],
         "c2_title": "Memory Mappings via mmap", "c2_bullets": ["• The 'mmap' system call allows processes to map file contents directly into their virtual memory.", "• This enables high-performance disk access: file read/write operations occur directly in RAM.", "• Permissions: page tables enforce read, write, and execute permissions on memory ranges."]},
         
        # Slide 5: COW Optimization
        {"type": "content2", "title": "The Copy-On-Write (COW) Design",
         "c1_title": "The RAM Optimization Strategy", "c1_bullets": ["• When a process requests a private file mapping using the 'MAP_PRIVATE' flag, memory is conserved.", "• Rather than duplication, the process's virtual memory maps to the same shared physical pages as the source file."],
         "c2_title": "Triggering Page Replication", "c2_bullets": ["• If a process attempts a write operation, the MMU triggers a page fault because the mapping is write-protected.", "• The kernel allocates a new physical frame, duplicates the source content, updates the process's page table to this new copy, and executes the write."]},
         
        # Slide 6: Vulnerability Mechanics
        {"type": "content2", "title": "The Vulnerability: COW Broken",
         "c1_title": "The Non-Atomic Page Fault Handler", "c1_bullets": ["• The virtual memory subsystem performs Copy-on-Write page allocations as a multi-step routine.", "• Multi-Step Sequence: (1) Check page permissions, (2) Trigger write page fault, (3) Allocate private page copy, (4) Update page table.", "• Problem: This multi-step process is non-atomic."],
         "c2_title": "Triggering the Race Condition", "c2_bullets": ["• While Thread B is executing the multi-step COW handler, Thread A repeatedly discards the memory map.", "• The kernel yields CPU control mid-operation, losing track of the COW allocation state.", "• Result: The write call writes directly to the original read-only backing cache."]},
         
        # Slide 7: Exploit Prerequisites
        {"type": "content2", "title": "Key System Calls & Components",
         "c1_title": "mmap and madvise Tools", "c1_bullets": ["• 'mmap' with 'MAP_PRIVATE' maps the target read-only root file into process virtual memory.", "• 'madvise' with 'MADV_DONTNEED' advises the kernel that the mapped page range is no longer needed.", "• The kernel responds by immediately clearing virtual page table links."],
         "c2_title": "The /proc/self/mem Backdoor", "c2_bullets": ["• Direct pointer writes (e.g. *ptr = 'x') fail instantly with MMU hardware segmentation faults.", "• Exploit: Bypasses the MMU by writing through '/proc/self/mem' (process virtual memory file).", "• This routes the write call directly through kernel memory routines."]},
         
        # Slide 8: Step-by-Step Exploit Flow
        {"type": "content2", "title": "Step-by-Step Race Condition Flow",
         "c1_title": "Thread A (The Discarder)", "c1_bullets": ["• Executes in a high-speed, infinite loop.", "• System Call: 'madvise(map, size, MADV_DONTNEED)'.", "• Action: Repeatedly clears the physical page table pointers for our mapped virtual memory address."],
         "c2_title": "Thread B (The Writer)", "c2_bullets": ["• Runs concurrently in a parallel, high-speed loop.", "• Action: Performs writes to `/proc/self/mem` targeting our mapped virtual address.", "• Timing Collision: The write permission check succeeds, the page is discarded by Thread A, and the write executes directly to read-only memory."]},
         
        # Slide 9: Container Security Insights
        {"type": "content2", "title": "Dirty COW inside Docker Containers",
         "c1_title": "The Shared Kernel Flaw", "c1_bullets": ["• Docker containers share the physical host system's Linux kernel directly, relying on namespaces for isolation.", "• If the host's underlying kernel is unpatched, every container running on it is exposed."],
         "c2_title": "Bypassing Namespace Isolation", "c2_bullets": ["• Running 'dirtyc0w' inside a container bypasses Docker's read-only mounts.", "• The exploit modifies the shared kernel page cache directly in physical RAM.", "• Changes propagate to the host file system, breaking the container boundary."]},
         
        # Slide 10: Privilege Escalation Escape
        {"type": "content2", "title": "Privilege Escalation to Container Escape",
         "c1_title": "Step 1: Container Root Elevation", "c1_bullets": ["• Attacker secures a shell as a low-privileged container user ('victim').", "• Running 'dirtyc0w' overwrites a container root configuration file, elevating the shell to container root."],
         "c2_title": "Step 2: Complete Host Breakout", "c2_bullets": ["• From container root, the attacker abuses administrative permissions (like CAP_SYS_ADMIN).", "• They target host sockets (like /var/run/docker.sock) to launch a privileged container.", "• They mount the host filesystem and establish persistent host root shells."]},
         
        # Slide 11: Exploit C Code Breakdown
        {"type": "content2", "title": "Exploit Implementation: C Code Breakdown",
         "c1_title": "Memory Mapping & Thread Spawning", "c1_bullets": ["• target opened as O_RDONLY; mapped with 'mmap(..., PROT_READ, MAP_PRIVATE, f, 0)'.", "• Spawns dual race threads using 'pthread_create'.", "• Thread 1: loops 'madvise(map, 100, MADV_DONTNEED)' 100,000 times."],
         "c2_title": "/proc/self/mem Writing Loop", "c2_bullets": ["• Thread 2: opens '/proc/self/mem' in O_RDWR mode.", "• Loops 100,000 times: calls 'lseek' to map address, then calls 'write' to inject payload.", "• Parallel core execution maximizes the probability of hitting the race window."]},
         
        # Slide 12: Hands-on Lab Environment
        {"type": "content2", "title": "Hands-on Lab: Container Environment",
         "c1_title": "Laboratory Sandbox Configuration", "c1_bullets": ["• Base OS: Ubuntu Docker container mapped to a vulnerable host kernel.", "• Unprivileged environment: runs under custom 'USER victim' shell context.", "• Target asset: read-only target file '/var/secret/target.txt' owned by root."],
         "c2_title": "CI/CD Pipeline Integrity", "c2_bullets": ["• The laboratory includes automated build scripts to compile the exploit binary safely.", "• Teaches developers and administrators how runtime boundaries fail in unpatched systems."]},
         
        # Slide 13: Executing the Lab
        {"type": "content2", "title": "Executing the Lab & Outcome Analysis",
         "c1_title": "Exploit Execution Flow", "c1_bullets": ["• Low-privileged 'victim' shell executes the compiled 'dirtyc0w' binary.", "• Shell command: './dirtyc0w /var/secret/target.txt PAYLOAD_STRING'.", "• The race threads execute concurrently, outputting completion counts."],
         "c2_title": "Vulnerability Proof", "c2_bullets": ["• Read-only file read check: 'cat /var/secret/target.txt'.", "• Outcome: The file content is overwritten by the payload.", "• Direct proof that kernel-level page protection has been bypassed successfully."]},
         
        # Slide 14: Real-world Exploitation cases
        {"type": "content2", "title": "Real-world Exploitation & Notable Cases",
         "c1_title": "Massive Cloud Infrastructure Impact", "c1_bullets": ["• Millions of Linux cloud servers, virtualization nodes, and Android mobile devices were left vulnerable.", "• Multi-tenant hypervisors were highly exposed: attackers could read/write physical memory of neighboring tenants."],
         "c2_title": "Wild Exploitation", "c2_bullets": ["• In-the-wild threat telemetry detected active mobile exploits bypassing OS permissions.", "• Emphasizes that host patching is the most critical foundation of container security."]},
         
        # Slide 15: Defense & Hardening Strategies
        {"type": "content2", "title": "Defense & Hardening Strategies",
         "c1_title": "Runtime Remediation", "c1_bullets": ["• Capability Dropping: Drop administrative capabilities: run containers with '--cap-drop=ALL'.", "• Read-only Runtimes: Mark root filesystems as read-only using '--read-only' to block exploit downloads.", "• Non-root users: Enforce UID limits."],
         "c2_title": "System-level Protections", "c2_bullets": ["• Kernel Patching: Apply updates instantly to host operating systems.", "• Linux Security Modules (LSMs): deploy AppArmor/SELinux profiles to restrict memory access calls."]},
         
        # Slide 16: Conclusion & Core Lessons
        {"type": "content2", "title": "Conclusion & Core Lessons",
         "c1_title": "Main Takeaways", "c1_bullets": ["• Containers are not sandboxes: direct kernel sharing represents an implicit trust relationship with the host.", "• Defense-in-depth is mandatory: runtime lockdowns must align with prompt kernel security patching."],
         "c2_title": "Securing the Pipeline", "c2_bullets": ["• supply chain provenance: use Cosign cryptographic signing to prevent unauthorized container execution.", "• Hexa Force Blueprint: delivers a hardened runtime paired with OIDC verified builds."]}
    ]
    
    # ----------------------------------------------------
    # DECK 2 Speaker Notes Content Mapped Slide-by-Slide
    # ----------------------------------------------------
    d2_notes_text = [
        # Slide 1
        "=== SPOKEN TALK TRACK ===\n"
        "\"Welcome to Part 2 of our briefing: The Dirty COW Technical Deep Dive. In this session, we will demystify the "
        "precise mechanisms behind CVE-2016-5195, commonly known as Dirty COW. This kernel exploit bypassed decades of "
        "security policies in Linux systems by weaponizing a race condition inside the virtual memory management subsystem. "
        "We will explore how memory pages are mapped, how Copy-on-Write page cache mechanisms work, and how we can force "
        "the kernel to write user data directly to read-only root files. Let us dive into the exploit's background.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Introduce the briefing as a deep, academically rigorous analysis of a legendary kernel vulnerability.\n"
        "• Frame the presentation around core operating system structures: virtual memory, page fault handlers, and thread scheduling.\n"
        "• Clearly declare the scope: it covers the CVE profile, exploitation theory, code walkthrough, and kernel patch analysis.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why is it called \"Dirty COW\"?\n"
        "Answer: \"COW\" refers to the \"Copy-on-Write\" mechanism. \"Dirty\" refers to marking a page in physical memory as modified in the page tables, indicating that its changes must be written back to the disk.",
        
        # Slide 2
        "=== SPOKEN TALK TRACK ===\n"
        "\"The briefing today is structured to provide both foundational theory and a hands-on lab demonstration. "
        "We will start with memory management basics, move through the COW design and vulnerability analysis, "
        "perform a deep code review of dirtyc0w.c, showcase our hands-on Docker laboratory environment, "
        "analyze the real-world impact, and outline concrete remediation strategies.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Lay out the technical roadmap for the deep dive to establish structured, expert-level delivery.\n"
        "• Note that the briefing bridges theory (operating system architecture) with practice (compiled exploits in containers).\n"
        "• Emphasize that understanding the vulnerability code is key to understanding modern mitigation strategies.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why is a deep dive on a 2016 vulnerability still highly valuable today?\n"
        "Answer: Dirty COW is the canonical textbook example of a virtual memory race condition exploit. By understanding it, security engineers gain deep insights into how kernel memory mapping works and why container environments (which share the host kernel) are highly vulnerable to shared resource exploits.",
        
        # Slide 3
        "=== SPOKEN TALK TRACK ===\n"
        "\"To contextualize this threat, let's examine the profile of CVE-2016-5195. Discovered in 2016, this "
        "vulnerability carries a High CVSS score. It affected almost every Linux kernel version prior to October 2016. "
        "The beauty—and danger—of Dirty COW is that it requires absolutely no special system privileges to execute. Any "
        "local user, even an unprivileged process running inside a Docker container, can trigger this exploit to bypass "
        "file permissions and write root-owned configurations. Let us look at the memory architecture that made this possible.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Cite precise vulnerability metadata: CVE-2016-5195, High severity CVSS rating, and pre-2016 kernel versions.\n"
        "• Highlight that this is a local privilege escalation (LPE) vector, requiring zero administrative starting permissions.\n"
        "• Set up the core threat: it completely compromises local system file integrity.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Can Dirty COW be executed over the network without local access?\n"
        "Answer: No, it is strictly a Local Privilege Escalation (LPE) vulnerability. The attacker must already have local shell access or code execution capabilities on the system to compile and run the exploit binary.",
        
        # Slide 4
        "=== SPOKEN TALK TRACK ===\n"
        "\"Let's review Linux virtual memory management. The operating system uses virtual memory addresses to isolate "
        "processes from each other and from physical hardware. Physical RAM is divided into fixed-size units called page "
        "frames, while virtual memory is managed in pages. The CPU's Memory Management Unit (MMU) and page tables translate "
        "virtual addresses to physical RAM addresses. Processes map files into their virtual memory using mmap, which allows "
        "reading or writing file data directly through RAM operations.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Explain the role of the MMU (Memory Management Unit) in performing address translations.\n"
        "• Define key terms: page frames (physical RAM), pages (virtual addresses), and page table entries.\n"
        "• Establish how `mmap` optimizes file reading by loading disk data directly into memory maps.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What role does the MMU play in memory protection?\n"
        "Answer: The MMU translates addresses and enforces page-level access permissions (e.g. read-only, read-write, executable) defined in the page table entries, throwing a segmentation fault if an unauthorized write is attempted.",
        
        # Slide 5
        "=== SPOKEN TALK TRACK ===\n"
        "\"Copy-on-Write is a crucial memory optimization. When a process maps a file privately using the MAP_PRIVATE flag, "
        "the kernel does not create a duplicate copy of the file in RAM. Instead, it maps the process's virtual memory "
        "to the same physical page cache frame as the underlying file to save memory. When the process attempts to write "
        "to this page, the kernel detects the write, allocates a new physical page frame, copies the original page content, "
        "updates the process's page table to point to this new private copy, and executes the write. This is the Copy-on-Write process.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Explain Copy-on-Write (COW) as an elegant RAM optimization technique used universally in modern operating systems.\n"
        "• Define `MAP_PRIVATE`: it guarantees that write modifications are made to a private memory space, never touching the underlying file.\n"
        "• Describe the transition from shared physical pages to private copy allocation during a write operation.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How does MAP_PRIVATE protect the original backing file?\n"
        "Answer: MAP_PRIVATE guarantees that all write operations are executed against a private local copy allocated in physical RAM, ensuring that the original backing file on disk is never modified.",
        
        # Slide 6
        "=== SPOKEN TALK TRACK ===\n"
        "\"The core vulnerability is a non-atomic operation in the page fault handler. When a process writes to a private "
        "read-only page, the kernel follows a multi-step routine: (1) Find the physical page, (2) Verify write permissions "
        "(which fails, triggering a page fault), (3) Perform the Copy-on-Write allocation, and (4) Update page table entries. "
        "Because this process is not atomic, the kernel must yield or check permissions at multiple boundaries. If page "
        "table entries are cleared concurrently during these checks, the kernel gets confused, allowing a concurrent write "
        "to write directly to the original read-only page before the COW is completed.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Define \"non-atomic\": a sequence of operations that can be interrupted or preempted by other CPU instructions.\n"
        "• Describe the precise moment the vulnerability occurs: the gap between checking the write permission and updating the page table entry.\n"
        "• Explain how concurrent threads can corrupt this gap, causing the kernel to fall back to writing to the read-only backing page.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What does non-atomic mean in this context?\n"
        "Answer: It means the operation is broken into multiple individual CPU instructions. In a multi-threaded system, other threads can execute instructions in between these steps, altering the memory state before the operation is completed.",
        
        # Slide 7
        "=== SPOKEN TALK TRACK ===\n"
        "\"The exploit utilizes two key system calls: (1) mmap with the MAP_PRIVATE flag to map the target read-only "
        "file into virtual memory, and (2) madvise with the MADV_DONTNEED flag. The madvise system call tells the kernel "
        "that the process no longer needs the specified memory range. The kernel responds by discarding the physical page "
        "mappings, forcing the next read to load the data directly from the original backing file. Additionally, the "
        "exploit writes directly to `/proc/self/mem` to bypass standard MMU write restrictions.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Detail the roles of `mmap(MAP_PRIVATE)` and `madvise(MADV_DONTNEED)`.\n"
        "• Explain that `MADV_DONTNEED` forces the kernel to instantly invalidate the physical page tables mapping.\n"
        "• Introduce `/proc/self/mem` as the direct virtual process memory interface used to trigger the kernel-level write routing.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why is writing to `/proc/self/mem` used instead of direct pointer dereferencing?\n"
        "Answer: Direct pointer writes are blocked by the CPU's MMU immediately, raising a segmentation fault. Writing through the `/proc/self/mem` interface routes the write through the kernel's virtual memory system, forcing the kernel itself to handle the write and execute the vulnerable page fault routine.",
        
        # Slide 8
        "=== SPOKEN TALK TRACK ===\n"
        "\"The exploit sequence runs two concurrent threads. Thread A runs in a tight loop calling `madvise(MADV_DONTNEED)` "
        "on our mapped memory address, repeatedly throwing away any newly allocated COW pages. Thread B runs in a parallel "
        "loop writing our payload to `/proc/self/mem` at the same virtual address. As they execute millions of times per "
        "second across CPU cores, a race condition occurs. The kernel completes the write permission check but gets "
        "preempted by madvise discarding the page. The write then executes directly against the original, read-only physical "
        "backing page, overwriting the file on disk.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Detail the concurrent execution flow: Thread A (`madvise`) and Thread B (`/proc/self/mem` write) racing at high speed.\n"
        "• Explain the physical collision: the write is directed to the read-only backing cache instead of the private copy.\n"
        "• Emphasize that the exploit requires multiple iterations in a loop to guarantee hitting the microsecond race window.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How many CPU cores are required to exploit the race condition?\n"
        "Answer: While it is theoretically possible on a single core via context switching, it is highly efficient on multi-core systems where Thread A and Thread B can run simultaneously on separate physical cores, maximizing the probability of racing the kernel's instructions.",
        
        # Slide 9
        "=== SPOKEN TALK TRACK ===\n"
        "\"Dirty COW is particularly catastrophic in Docker container environments. Because Docker containers share the "
        "host's Linux kernel, any container—even one running as a low-privileged user—can run `dirtyc0w` to exploit the "
        "shared kernel. The exploit can bypass container isolation barriers entirely. An attacker inside a container "
        "can overwrite a read-only system file on the host's kernel or overwrite physical files shared between containers, "
        "collapsing the foundational promise of container isolation.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Explain why shared kernels represent the central security vulnerability of containerization.\n"
        "• Discuss the collapse of container isolation boundaries when a container process can overwrite host-level system files.\n"
        "• Reiterate that namespaces and cgroups do not prevent physical page cache corruption because both container and host share the same kernel cache.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How does Dirty COW bypass Docker's read-only file mappings?\n"
        "Answer: Docker's read-only mappings are enforced at the mount and namespace levels. However, because the exploit operates at the shared kernel virtual memory page level, it writes directly to the shared physical page cache. The kernel executes the physical write, rendering namespace and mount restrictions completely ineffective.",
        
        # Slide 10
        "=== SPOKEN TALK TRACK ===\n"
        "\"Let's trace the threat pathway from privilege escalation to complete container escape. First, the attacker "
        "exploits an application vulnerability to gain a shell as a low-privileged container user. Second, they run "
        "Dirty COW to overwrite a root-owned file (like `/etc/passwd` or `/var/secret/target.txt`), gaining root "
        "privileges inside the container. Third, they leverage container root to access administrative capabilities "
        "or mounts (like the exposed `/var/run/docker.sock`), allowing them to launch new privileged containers, mount "
        "the physical host filesystem, and achieve complete host escape.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Illustrate the escalate-and-escape pathway step-by-step.\n"
        "• Connect Dirty COW (privilege escalation) with local boundary breakages (capabilities and exposed sockets).\n"
        "• Emphasize that a container compromise almost always leads to a host compromise if runtime capabilities are left unchecked.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Is container escape possible without local privilege escalation?\n"
        "Answer: Yes, if the container is already run with high privileges (like '--privileged' or mounting the host filesystem directly). However, if the container is properly locked down, privilege escalation (like Dirty COW) is the necessary first step to unlock the administrative APIs required for escape.",
        
        # Slide 11
        "=== SPOKEN TALK TRACK ===\n"
        "\"Our C implementation `dirtyc0w.c` is elegant and compact. We open the target file, retrieve its size, and map "
        "it read-only but privately using `mmap`. We then spin up two parallel threads using `pthread_create`. Thread 1 "
        "runs `madviseThread` calling `madvise(map, 100, MADV_DONTNEED)` in a loop of 100,000 iterations. Thread 2 runs "
        "`procSelfMemThread` opening `/proc/self/mem`, seeking to the mapped address, and writing the payload in a loop of "
        "100,000 iterations. The parallel execution guarantees that the race condition will be triggered within seconds.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Walk through the core functions of `dirtyc0w.c`: `mmap`, `pthread_create`, `madviseThread`, and `procSelfMemThread`.\n"
        "• Explain the seek operation `lseek(f, (uintptr_t)map, SEEK_SET)` as the mapping pointer address offset aligner.\n"
        "• Detail the iteration scale (100,000 runs) that ensures the thread collision occurs quickly.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What happens to the target file if its size is smaller than the payload?\n"
        "Answer: The exploit can only overwrite data up to the original size of the mapped file. Attempting to write beyond the file boundary will fail or result in a bus error, which is why the exploit is typically used to overwrite configuration files (like `/etc/passwd`) where the text can be swapped in-place.",
        
        # Slide 12
        "=== SPOKEN TALK TRACK ===\n"
        "\"To demonstrate this exploit in a controlled environment, Hexa Force built a complete Docker-based laboratory. "
        "The laboratory uses a custom Dockerfile that copies `dirtyc0w.c` and a shell demo script `run_demo.sh` into an "
        "unprivileged container running under the user `victim`. A read-only file `/var/secret/target.txt` is created "
        "inside the container as the target. The environment is hosted on an unpatched Linux kernel to allow students "
        "and security engineers to safely observe the exploit in real-time.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Describe the laboratory sandbox components: unprivileged container environment, `victim` user, and read-only target file.\n"
        "• Explain that the container is run on an unpatched Linux kernel within a nested hypervisor or isolated virtual machine.\n"
        "• Emphasize that the lab design replicates real-world microservice exploitation in a safe, controlled manner.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How do you prevent students from accidentally compromising the host in this lab?\n"
        "Answer: The lab is run in an isolated virtual machine or sandbox environment. Although the container can theoretically exploit the kernel to escape, running the entire demonstration inside a dedicated test VM guarantees that any host compromise is contained within the sandbox.",
        
        # Slide 13
        "=== SPOKEN TALK TRACK ===\n"
        "\"Let's examine the outcome of executing the lab. When we run `./run_demo.sh`, the script compiles `dirtyc0w.c` "
        "using gcc and executes the binary against the read-only file `/var/secret/target.txt` with a replacement payload. "
        "The threads spin up, complete their 100,000 iterations, and print their completion. When we read `/var/secret/target.txt` "
        "immediately after, we find that the read-only file has been successfully overwritten, proving the exploit's success in real-time.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Outline the real-time execution steps: gcc compilation, target mapping, multi-threaded racing, and outcome reading.\n"
        "• Emphasize the immediate, visual feedback of the overwrite to prove exploit success.\n"
        "• Position the lab as a concrete proof of the kernel's virtual memory vulnerability under high-speed thread competition.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why does the exploit sometimes require multiple executions to succeed?\n"
        "Answer: Because it relies on a timing race condition. Depending on CPU scheduling, system load, and thread execution alignment, the race window may not be hit immediately. However, running the loops 100,000 times almost always guarantees success within a single run.",
        
        # Slide 14
        "=== SPOKEN TALK TRACK ===\n"
        "\"The real-world impact of Dirty COW was massive. It affected millions of Linux servers, Android devices, "
        "and cloud microservice environments worldwide. In multi-tenant environments—where multiple customers share the "
        "same physical server kernel—Dirty COW allowed malicious tenants to bypass isolation and access other tenants' data. "
        "It serves as a stark reminder of why prompt kernel patching is a fundamental requirement of cloud infrastructure security.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Discuss the global scope of CVE-2016-5195: millions of servers, Android mobile devices, and virtualization nodes compromised.\n"
        "• Explain the risk to multi-tenant cloud hosting: unprivileged users reading or writing adjacent customer memory segments.\n"
        "• Reiterate that kernel level patching is not an optional operation, but a foundational cloud security mandate.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Was Dirty COW ever exploited in the wild?\n"
        "Answer: Yes. Security researchers and threat intelligence firms detected in-the-wild exploits using Dirty COW to root Android devices, bypass mobile security controls, and compromise multi-tenant web hosting servers before patches were universally applied.",
        
        # Slide 15
        "=== SPOKEN TALK TRACK ===\n"
        "\"Remediating Dirty COW requires a multi-layered security strategy. The primary defense is applying kernel updates "
        "to patch the vulnerability. Beyond kernel patching, we must implement Defense-in-Depth. In containers, we should "
        "strictly drop administrative capabilities (like CAP_SYS_ADMIN), run processes as non-root users, utilize read-only "
        "root filesystems, and employ AppArmor or SELinux profiles to restrict system calls like madvise or write access "
        "to `/proc/self/mem`.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Position the core solution: apply the Linux kernel security patch that fixes Copy-on-Write page table routines.\n"
        "• Walk through the supplemental Defense-in-Depth layers: capability dropping, non-root runtimes, and read-only filesystems.\n"
        "• Explain how LSMs (Linux Security Modules) like AppArmor or SELinux can block unauthorized syscalls like `madvise`.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How can AppArmor or SELinux mitigate Dirty COW if the kernel is unpatched?\n"
        "Answer: AppArmor or SELinux can block the process from writing to `/proc/self/mem` or restrict the execution of specific system calls, neutralising the exploit's primary vector even if the underlying kernel vulnerability is still present.",
        
        # Slide 16
        "=== SPOKEN TALK TRACK ===\n"
        "\"In conclusion, the Dirty COW technical briefing highlights the deep architectural risks of shared-kernel "
        "virtualization. While containerization offers unparalleled efficiency, it demands rigorous host-level patching, "
        "strict runtime isolation rules, and continuous security automation. The Hexa Force DevSecOps architecture "
        "successfully mitigates these risks, delivering a robust blueprint for modern enterprise security. Thank you, "
        "and we are open to any questions.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Deliver a strong technical summary linking shared-kernel realities to continuous security patching.\n"
        "• Reiterate that modern microservices demand runtime sandboxing, supply chain integrity, and regular kernel audits.\n"
        "• Open the floor to final technical questions with absolute poise and authority.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the most important lesson from the Dirty COW briefing?\n"
        "Answer: We must never treat containers as full security boundaries by default. Container security requires a continuous DevSecOps pipeline combining host kernel patching, least-privilege runtimes, and supply chain integrity validation to guarantee absolute system security."
    ]
    
    # Render Deck 2 Slide Designs
    blank_layout = prs.slide_layouts[6]
    
    for i, slide_data in enumerate(slides_d2):
        slide = prs.slides.add_slide(blank_layout)
        apply_dark_bg(slide)
        
        if slide_data["type"] == "title":
            add_title_slide(slide, slide_data["title"], slide_data["subtitle"], slide_data["cat"])
        else:
            add_header(slide, slide_data["title"], category="DIRTY COW TECHNICAL ASSESSMENTS")
            
            # Draw Two Columns
            col_w = Inches(5.6)
            col_h = Inches(4.8)
            top_y = Inches(1.8)
            
            # Left Card
            draw_card(slide, Inches(0.75), top_y, col_w, col_h)
            add_card_text(slide, Inches(0.75), top_y, col_w, col_h, slide_data["c1_title"], slide_data["c1_bullets"], title_color=C_PRIMARY)
            
            # Right Card
            border_c = C_GREEN if "Mitigation" in slide_data["title"] or "Remediation" in slide_data["title"] or "Defence" in slide_data["title"] or "Lessons" in slide_data["title"] else None
            draw_card(slide, Inches(6.8), top_y, col_w, col_h, border_color=border_c)
            add_card_text(slide, Inches(6.8), top_y, col_w, col_h, slide_data["c2_title"], slide_data["c2_bullets"], title_color=C_GREEN if border_c else C_PRIMARY)
            
        # Inject Speaker Notes
        notes_slide = slide.notes_slide
        tf_notes = notes_slide.notes_text_frame
        tf_notes.text = d2_notes_text[i]
        
    print("[+] Generated Deck 2 Slides programmatically.")

def main():
    print("[*] Rebuilding presentations with ultimate premium layout...")
    
    # Rebuild Deck 1
    prs1 = Presentation()
    prs1.slide_width = Inches(13.333)
    prs1.slide_height = Inches(7.5)
    build_deck_1(prs1)
    prs1.save("Hexa_Force_Secure_Lab_Presentation.pptx")
    print("[+] Saved Hexa_Force_Secure_Lab_Presentation.pptx successfully.")
    
    # Rebuild Deck 2
    prs2 = Presentation()
    prs2.slide_width = Inches(13.333)
    prs2.slide_height = Inches(7.5)
    build_deck_2(prs2)
    prs2.save("Dirty_COW_Technical_Briefing.pptx")
    print("[+] Saved Dirty_COW_Technical_Briefing.pptx successfully.")
    
    print("[+] ALL SLIDE DECKS COMPLETED!")

if __name__ == "__main__":
    main()
