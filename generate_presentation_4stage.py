import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 Widescreen standard
prs.slide_height = Inches(7.5)

# Color Scheme Definition
COLOR_BG = RGBColor(15, 23, 42)       # Slate 900 #0F172A
COLOR_WHITE = RGBColor(255, 255, 255) # Pure White
COLOR_SILVER = RGBColor(148, 163, 184)# Slate 400 #94A3B8
COLOR_CARD = RGBColor(30, 41, 59)     # Slate 800 #1E293B

# Accent colors for pillars
COLOR_RED = RGBColor(239, 68, 68)     # Accent Crimson (Attack) #EF4444
COLOR_CYAN = RGBColor(6, 182, 212)    # Accent Teal (Diagnostics) #06B6D4
COLOR_GREEN = RGBColor(34, 197, 94)   # Accent Green (Mitigation) #22C55E
COLOR_YELLOW = RGBColor(234, 179, 8)  # Accent Yellow #EAB308

# Helper function to set slide background
def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

# Helper function to create title section on content slides
def add_slide_header(slide, title_text, category="HEXA FORCE CONTAINER SECURITY"):
    # Add a thin cyan category label
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Trebuchet MS"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_CYAN
    
    # Add the main title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

# Helper to create bullet points in a modern look
def add_bullet_points(slide, left, top, width, height, points):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    for i, pt in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(10)
        
        # Check special prefixes
        if pt.startswith("  - "):
            p.text = "   •  " + pt[4:]
            p.font.name = "Calibri"
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_SILVER
        elif pt.startswith("[CRITICAL]"):
            p.text = pt.replace("[CRITICAL]", "🔴 ")
            p.font.name = "Calibri"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = COLOR_RED
        elif pt.startswith("[DIAGNOSTIC]"):
            p.text = pt.replace("[DIAGNOSTIC]", "🔵 ")
            p.font.name = "Calibri"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
        elif pt.startswith("[DEFENDED]"):
            p.text = pt.replace("[DEFENDED]", "🟢 ")
            p.font.name = "Calibri"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = COLOR_GREEN
        elif pt.startswith("[HIGHLIGHT]"):
            p.text = pt.replace("[HIGHLIGHT]", "⚡ ")
            p.font.name = "Calibri"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
        else:
            p.text = "•  " + pt
            p.font.name = "Calibri"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            
    return box

# Helper to draw a stylized code box
def add_code_box(slide, left, top, width, height, lines):
    # Add background card for code
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = COLOR_CYAN
    card.line.width = Pt(1)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(2)
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.LEFT
        
        # Simple color syntax highlights
        if "docker run" in line or "docker build" in line or "cosign sign" in line:
            p.font.color.rgb = COLOR_CYAN
        elif "--stage-" in line or "IMAGE_TAG" in line or "cosign-installer" in line:
            p.font.color.rgb = COLOR_GREEN
        elif "victim" in line or "github.repository" in line:
            p.font.color.rgb = COLOR_YELLOW
            
    return card

# Helper to draw visual info cards
def add_info_card(slide, left, top, width, height, title, subtitle, content, border_color):
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = title.upper()
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.space_after = Pt(4)
    
    # Subtitle
    if subtitle:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = border_color
        p_sub.space_after = Pt(8)
        
    # Content
    for text_line in content:
        p_cont = tf.add_paragraph()
        p_cont.text = "• " + text_line
        p_cont.font.name = "Calibri"
        p_cont.font.size = Pt(12)
        p_cont.font.color.rgb = COLOR_SILVER
        p_cont.space_after = Pt(4)

# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)

title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "HEXA FORCE CONTAINER LAB"
p1.font.name = "Trebuchet MS"
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = COLOR_CYAN
p1.space_after = Pt(4)

p2 = tf.add_paragraph()
p2.text = "4-Stage DevSecOps Attack & Mitigation Laboratory"
p2.font.name = "Trebuchet MS"
p2.font.size = Pt(26)
p2.font.color.rgb = COLOR_WHITE
p2.space_after = Pt(14)

p3 = tf.add_paragraph()
p3.text = "Presented by Hexa Force | Cryptographically Signed Container Architectures, OIDC Trust, and DevSecOps Containment Rules"
p3.font.name = "Calibri"
p3.font.size = Pt(13)
p3.font.italic = True
p3.font.color.rgb = COLOR_SILVER

# Metadata bottom
meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8))
tf_meta = meta_box.text_frame
p_meta = tf_meta.paragraphs[0]
p_meta.text = "PIPELINE ORCHESTRATION & LAB MANUAL  |  HEXA FORCE INTEGRATION BRIEFING"
p_meta.font.name = "Trebuchet MS"
p_meta.font.size = Pt(11)
p_meta.font.bold = True
p_meta.font.color.rgb = COLOR_SILVER

# ==========================================
# SLIDE 2: Agenda
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Briefing & Laboratory Agenda", "EXECUTIVE PRESENTATION GUIDE")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), [
    "01. The 4 Security Pillars Overview",
    "  - Key definitions, threat structures, and container security scope.",
    "02. Stage 1: Host Kernel Isolation",
    "  - Dirty COW (CVE-2016-5195) C exploit simulation and outcomes.",
    "03. Stage 2: Process & Capabilities Isolation",
    "  - Shared PID namespace boundaries and evaluation of CAP_SYS_PTRACE.",
    "04. Stage 3: Daemon & API Security",
    "  - Mitigating exposed docker.sock socket API mounts."
])

add_bullet_points(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), [
    "05. Stage 4: Volume & Filesystem Security",
    "  - Writable host mount points and path traversal cron injection.",
    "06. Cryptographic Container Signing",
    "  - Keyless Cosign OIDC authentication and trust chain verification.",
    "07. GitHub Actions CI/CD Integration",
    "  - Case-compliant image compilation, signing, and pipeline execution.",
    "08. Hexa Force Conclusion",
    "  - Key lessons, takeaways, and final architecture summaries."
])

# ==========================================
# SLIDE 3: The 4 Security Pillars
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "The 4 Security Pillars of Containers", "LABORATORY SCOPE")

card_w = Inches(2.7)
card_h = Inches(4.8)

add_info_card(
    slide, Inches(0.8), Inches(1.8), card_w, card_h,
    "1. Host Kernel", "Isolation Boundary",
    [
        "Containers share host OS kernel resources.",
        "Race condition exploits (e.g. Dirty COW) bypass all virtual restrictions.",
        "Demonstrates kernel memory isolation limits."
    ],
    COLOR_RED
)

add_info_card(
    slide, Inches(3.8), Inches(1.8), card_w, card_h,
    "2. Capabilities", "Process Boundaries",
    [
        "Process namespace mapping regulates visibility.",
        "Over-privileged debug access (CAP_SYS_PTRACE) allows host code injection.",
        "Demonstrates process capability containment."
    ],
    COLOR_CYAN
)

add_info_card(
    slide, Inches(6.8), Inches(1.8), card_w, card_h,
    "3. Engine API", "Daemon Protection",
    [
        "Mounting docker.sock grants host root access.",
        "Attacker controls engine and launches sibling mounts.",
        "Demonstrates socket boundary control."
    ],
    COLOR_YELLOW
)

add_info_card(
    slide, Inches(9.8), Inches(1.8), card_w, card_h,
    "4. Volume", "Filesystem Isolation",
    [
        "Writable mounts allow local file traversal.",
        "Exploit writes root task files to host cron.",
        "Demonstrates mount point security rule options."
    ],
    COLOR_GREEN
)

# ==========================================
# SLIDE 4: Stage 1 - Host Kernel Isolation
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Stage 1: Host Kernel Isolation (Dirty COW)", "STAGE 1 WORKFLOW")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Exploit Concept (CVE-2016-5195)",
    "  - A kernel-level race condition in Linux's Copy-on-Write (COW) system.",
    "  - Overwrites read-only cache pages using procfs self-memory loops.",
    "[CRITICAL] Exploit Execution",
    "  - Low-privilege container user compiles and runs dirtyc0w.c.",
    "  - Targets a read-only root file: /var/secret/target.txt.",
    "[DEFENDED] Automated Protection",
    "  - On modern kernels (like GitHub runners), page faults are atomic.",
    "  - The race condition fails safely (Exploit Prevented)."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Detection & Mitigation Blueprint", "Hexa Force Remediation",
    [
        "Detection: Static scan container kernels for CVE-2016-5195 signatures.",
        "Detection: Audit syscall sequences using seccomp/madvise logs.",
        "Mitigation: Keep the underlying host OS kernel fully patched.",
        "Mitigation: Implement secure sandboxed runtimes (gVisor) to isolate guest system calls."
    ],
    COLOR_RED
)

# ==========================================
# SLIDE 5: Stage 2 - Process & Capabilities Isolation
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Stage 2: Process & Capabilities Isolation", "STAGE 2 WORKFLOW")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Exploit Concept (CAP_SYS_PTRACE)",
    "  - Over-privileged capabilities allow reading and writing process memory.",
    "  - Combined with shared host PID namespace, it allows host code injection.",
    "[DIAGNOSTIC] Capability Diagnostics",
    "  - Evaluates active permissions using the 'capsh --print' utility.",
    "  - Checks PID visibility to detect shared namespace profiles.",
    "[DEFENDED] Namespace Isolation",
    "  - With standard isolated namespaces, host processes are invisible.",
    "  - Ensures container tasks are safely isolated."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Detection & Mitigation Blueprint", "Hexa Force Remediation",
    [
        "Detection: Scan configuration files for '--pid=host' namespace flags.",
        "Detection: Monitor running processes for active CAP_SYS_PTRACE capabilities.",
        "Mitigation: Drop all default capabilities: --cap-drop=ALL.",
        "Mitigation: Strictly avoid running containers in host PID namespaces."
    ],
    COLOR_CYAN
)

# ==========================================
# SLIDE 6: Stage 3 - Daemon & API Security
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Stage 3: Daemon & API Security (docker.sock)", "STAGE 3 WORKFLOW")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Exploit Concept (UNIX Socket Mount)",
    "  - The Docker daemon runs as root on the host system.",
    "  - Exposing docker.sock gives the container full daemon administrative access.",
    "[CRITICAL] Sibling Hijack Simulation",
    "  - Detects exposed sockets and constructs REST API payloads.",
    "  - Spawns sibling containers with host root mounted on guest filesystems.",
    "[DEFENDED] Socket Protection",
    "  - Verifies that without socket file exposure, API queries are blocked.",
    "  - Guest container remains securely separated."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Detection & Mitigation Blueprint", "Hexa Force Remediation",
    [
        "Detection: Scan container mount definitions for docker.sock references.",
        "Detection: Implement runtime filters (e.g. Falco) to log raw curl queries to sockets.",
        "Mitigation: Prohibit socket mounts in unprivileged workloads.",
        "Mitigation: Run Docker in Rootless Mode so the daemon has no host root power."
    ],
    COLOR_YELLOW
)

# ==========================================
# SLIDE 7: Stage 4 - Volume & Filesystem Security
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Stage 4: Volume & Filesystem Security", "STAGE 4 WORKFLOW")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Exploit Concept (Writable Mounts)",
    "  - Writable host directory mounts allow guess writes to local system files.",
    "  - Compromised container writes scheduled tasks to host cron folders.",
    "[CRITICAL] Cron Injection Attack",
    "  - Simulates attempts to write reverse shell scripts to /mnt/host/etc/cron.d.",
    "  - On success, the host runs the arbitrary script as root in 60s.",
    "[DEFENDED] Read-Only Constraints",
    "  - Mount is configured as read-only, returning 'Permission Denied' blocks.",
    "  - Verifies that read-only volumes perfectly contain directory traversals."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Detection & Mitigation Blueprint", "Hexa Force Remediation",
    [
        "Detection: Configure File Integrity Monitoring (FIM) on host /etc/cron* directories.",
        "Detection: Set up auditd rules to log write activity on host cron folders.",
        "Mitigation: Always mount host directories strictly as read-only (:ro).",
        "Mitigation: Avoid mounting root directories (/etc, /var/run, /root) into containers."
    ],
    COLOR_GREEN
)

# ==========================================
# SLIDE 8: Cryptographic Container Signing (Cosign & OIDC)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Cryptographic Image Signing (Cosign & OIDC)", "SUPPLY CHAIN INTEGRITY")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Secure Supply Chain",
    "  - Signatures protect against registry compromises, spoofing, and tagging attacks.",
    "  - Verifies that only authorized, scanned images are executed in clusters.",
    "[HIGHLIGHT] Keyless OIDC Authentication",
    "  - Connects GitHub's OpenID Connect (OIDC) identity with Sigstore's CA (Fulcio).",
    "  - Eliminates private key storage, rotation, and leakage risks entirely.",
    "Short-Lived Certificates",
    "  - Fulcio issues temporary 10-minute signing certificates based on runner JWTs.",
    "  - The signature is pushed to the registry and verified keylessly using public APIs."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "OIDC Cryptographic Trust Chain", "Hexa Force Keyless signing",
    [
        "1. Runner requests JWT from GitHub Actions identity provider.",
        "2. Presents JWT token to Sigstore's certificate authority.",
        "3. Sigstore verifies GHA identity and generates short-lived cert.",
        "4. Image is cryptographically signed and signature pushed to GHCR.",
        "5. Runtimes verify authenticity keylessly using GitHub public APIs."
    ],
    COLOR_CYAN
)

# ==========================================
# SLIDE 9: CI/CD Pipeline Integration
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Automated CI/CD Laboratory Pipeline", "PIPELINE ORCHESTRATION")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "DevSecOps Orchestration",
    "  - Fully integrated using GitHub Actions workflow automations.",
    "  - Builds, tags, logs in, signs, and executes lab stages on commit.",
    "[HIGHLIGHT] Lowercase Tag Compliance",
    "  - Enforces strict lowercase registry tags using shell string conversions.",
    "  - Resolves Docker build case-sensitivity requirements automatically.",
    "Fast Execution Loop",
    "  - Optimizations reduce execution duration to under 50 seconds.",
    "  - Provides instant DevSecOps regression security assessment feedback."
])

add_code_box(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), [
    "# Set lowercase repository name for Docker/GHCR",
    "- name: Set lowercase repository name",
    "  run: |",
    "    echo \"IMAGE_TAG=ghcr.io/${GITHUB_REPOSITORY,,}\" >> ${GITHUB_ENV}",
    "",
    "# Install Cosign (Sigstore)",
    "- name: Install Cosign",
    "  uses: sigstore/cosign-installer@v3.5.0",
    "",
    "# Build & Push Container",
    "- name: Build and Push Docker Image",
    "  uses: docker/build-push-action@v5",
    "  with:",
    "    push: true",
    "    tags: ${{ env.IMAGE_TAG }}:latest",
    "",
    "# Keyless Sign the Published Image",
    "- name: Sign the Published Image",
    "  run: cosign sign --yes \"${{ env.IMAGE_TAG }}:${{ github.sha }}\""
])

# ==========================================
# SLIDE 10: DevSecOps Hardening Rules
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "DevSecOps Hardening & Defense-in-Depth", "BEST PRACTICES GUIDE")

# Three horizontal cards
add_info_card(
    slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
    "1. Host Controls", "Kernel Security Level",
    [
        "Apply host kernel updates immediately to resolve CVE-2016-5195.",
        "Implement seccomp security rules to block madvise MADV_DONTNEED if not required.",
        "Employ sandbox container engines (gVisor) in untrusted public environments."
    ],
    COLOR_RED
)

add_info_card(
    slide, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8),
    "2. Container Configuration", "Least Privilege Isolation",
    [
        "Never employ --pid=host namespaces.",
        "Drop all default privileges: docker run --cap-drop=ALL.",
        "Run services as custom USER victim (unprivileged accounts)."
    ],
    COLOR_CYAN
)

add_info_card(
    slide, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8),
    "3. Storage & Storage APIs", "Access Controls",
    [
        "Enforce read-only volumes strictly using the :ro flag.",
        "Prohibit mounting /var/run/docker.sock UNIX sockets.",
        "Avoid mounting critical host configurations (/etc, /opt, /root)."
    ],
    COLOR_GREEN
)

# ==========================================
# SLIDE 11: Conclusion & Key Takeaways
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Conclusion: Hexa Force Lab Debriefing", "KEY TAKEAWAYS")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Containers Share Host OS Foundations",
    "  - Container containment is only as strong as the host kernel.",
    "  - Bypassing kernel barriers compromises all namespaces.",
    "Mitigations are Multi-layered",
    "  - Securing pipelines requires host updates, capabilities control,",
    "  - socket restrictions, and read-only storage settings.",
    "[DEFENDED] Successful Hexa Force Outcome",
    "  - The CI/CD pipeline verified that modern kernels and proper",
    "  - read-only volume controls successfully contain advanced escapes.",
    "Continuous Security Regression Testing",
    "  - Integrating automated lab checks in GHA stops regressions before deploy."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "The Golden Rule of DevSecOps", "Final Recommendation",
    [
        "Containers are NOT isolation virtual machines.",
        "Host OS updates are the baseline requirement for container cluster security.",
        "Drop capabilities, run as unprivileged users, mount volumes read-only, and audit docker.sock exposure constantly."
    ],
    COLOR_CYAN
)

# Save the presentation
output_path = "Hexa_Force_Secure_Lab_Presentation.pptx"
prs.save(output_path)
print(f"[+] Successfully generated widescreen 16:9 presentation at: {os.path.abspath(output_path)}")
