import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------
# Brand Theme Color Palette (IITJ Academic-Executive Dark Mode)
# ----------------------------------------------------
C_BG = RGBColor(15, 23, 42)        # Slate 900 (Dark Slate Background)
C_CARD = RGBColor(30, 41, 59)      # Slate 800 (Card Panel Background)
C_PRIMARY = RGBColor(59, 130, 246)  # Electric Blue (IITJ Slate Blue Accent)
C_GOLD = RGBColor(245, 158, 11)     # Amber Gold (IITJ Premium Gold Accent)
C_TEXT = RGBColor(241, 245, 249)   # Cool White (Body Text)
C_MUTED = RGBColor(148, 163, 184)  # Muted Gray (Subtitles/Sidebars)
C_GREEN = RGBColor(16, 185, 129)   # Emerald Green (Success Indicators)
C_RED = RGBColor(239, 68, 68)      # Crimson Red (Threats/Risks)
C_WHITE = RGBColor(255, 255, 255)  # Pure White

# ----------------------------------------------------
# Slide Layout Helper Functions
# ----------------------------------------------------
def apply_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG

def add_header(slide, title, category="HEXA FORCE CONTAINER SECURITY  |  IITJ STAFF BRIEFING"):
    # Category Tag
    cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(9)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_GOLD
    
    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(11.83), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE

def draw_card(slide, left, top, width, height, bg_color=C_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_card_text(slide, left, top, width, height, title, body_bullets, title_color=C_WHITE, font_size_body=13):
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Card Title
    if title:
        p_title = tf.paragraphs[0]
        p_title.text = title.upper()
        p_title.font.name = "Arial"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.space_after = Pt(10)
        
    for bullet in body_bullets:
        p = tf.add_paragraph() if tf.paragraphs[-1].text else tf.paragraphs[0]
        p.text = bullet
        p.font.name = "Arial"
        p.font.size = Pt(font_size_body)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(6)
        
        if bullet.strip().startswith("•") or bullet.strip().startswith("-"):
            p.level = 0
        else:
            p.level = 0

def add_title_slide(slide, main_title, subtitle, category="HEXA FORCE RESEARCH BRIEFING"):
    apply_dark_bg(slide)
    
    # Structural Accent Band
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(0.12), Inches(3.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_PRIMARY
    accent.line.fill.background()
    
    # Category
    cat_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.0), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_GOLD
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.0), Inches(1.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = main_title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.0), Inches(1.0))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = C_MUTED

# ----------------------------------------------------
# PPTX Generator Data for all 15 slides
# ----------------------------------------------------
def build_final_briefing():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Slide data mapping (Outline 1 - 15)
    slides_data = [
        # Slide 1: Cover
        {"type": "title", "title": "HEXA FORCE CONTAINER SECURITY LAB", "subtitle": "A 4-Stage DevSecOps Attack, Detection, and Defense Laboratory\nPrepared for IITJ Staff Briefing  |  May 2026", "cat": "IITJ Security Research Series"},
        
        # Slide 2: Executive Summary
        {"type": "content2", "title": "Executive Summary & Headline Outcomes",
         "c1_title": "Core System Achievements", "c1_bullets": ["• Rebuilt complete 4-stage container breakout exploit path, validating namespaces, capabilities, sockets, and mounts.", "• Implemented proactive system call containment using custom Seccomp security profiles to block 'madvise'.", "• Engineered cryptographically signed OIDC container supply chains."],
         "c2_title": "DevSecOps Metrics", "c2_bullets": ["• Exploit Containment Rate: 100% successfully mitigated under Seccomp runtime filters.", "• Image Authenticity Rate: 100% of images scanned by Trivy and keylessly signed by Cosign.", "• Zero false positives reported during live laboratory demonstrations."]},
         
        # Slide 3: Problem/Context
        {"type": "content2", "title": "Problem Statement & Threat Landscape",
         "c1_title": "Shared Host Kernel Realities", "c1_bullets": ["• Core architectural flaw: standard Docker containers share the host Linux kernel directly.", "• Any local privilege escalation (LPE) vulnerability in the host kernel (like Dirty COW) is shared by all containers.", "• Logical boundaries (namespaces) fall instantly to physical page cache corruption exploits."],
         "c2_title": "The Exploit Cascade", "c2_bullets": ["• An unprivileged web compromised shell is elevated to container root.", "• Attacker leverages internal capabilities (CAP_SYS_ADMIN) and exposed Unix domain sockets (/var/run/docker.sock).", "• Spawns privileged containers, mounts physical host filesystem, and takes host root control."]},
         
        # Slide 4: Objectives & KPIs
        {"type": "content2", "title": "Objectives & KPI Target Matrix",
         "c1_title": "Security KPI Targets", "c1_bullets": ["• Target 1: Reduce host escape likelihood to 0% under active runtime policies.", "• Target 2: Achieve 100% cryptographic supply chain provenance via OIDC.", "• Target 3: Automated image linting and vulnerability scanning (Trivy) on push."],
         "c2_title": "Baseline vs Actual Results", "c2_bullets": ["• Host Breakout Escape: Baseline 100% (Vulnerable) $\rightarrow$ Actual 0% (Hardened).", "• Cryptographic Provenance: Baseline 0% (Unsigned) $\rightarrow$ Actual 100% (Fulcio/Rekor).", "• Vulnerability Scanning: Baseline 0% (Blind Trust) $\rightarrow$ Actual 100% (Trivy scan)."]},
         
        # Slide 5: Solution Overview
        {"type": "content2", "title": "DevSecOps Solution Architecture",
         "c1_title": "Continuous Pipeline Security", "c1_bullets": ["• Local Build Gate: Compiles container and triggers automated static checks.", "• AquaSec Trivy Scan: Evaluates vulnerabilities before distribution.", "• Cosign OIDC Exchange: GHA exchanges identity tokens for Fulcio CA certs, logging to Rekor."],
         "c2_title": "System Runtime Hardening", "c2_bullets": ["• Syscall Sandboxing: Custom Seccomp profile blocks the vulnerable 'madvise' system call.", "• Least Privilege Enforcement: Stripping capabilities and enforcing non-root executions.", "• Volume Governance: Mounting external directories as read-only."]},
         
        # Slide 6: Scope & Deliverables
        {"type": "content2", "title": "Project Scope & Deliverables Checklist",
         "c1_title": "Laboratory & Code Assets", "c1_bullets": ["• Hardened Dockerfile: Sets up non-root sandbox and read-only targets.", "• C Exploit Source (dirtyc0w.c): Multi-threaded memory mapped race condition demo code.", "• Custom Seccomp Profile (hexaforce-seccomp.json): blocks 'madvise' system calls."],
         "c2_title": "Orchestrators & Documentation", "c2_bullets": ["• Lab Runner (run_demo.sh): Sequentially models and executes all 4 stages.", "• DevSecOps Pipeline (.github/workflows/dirtycow-lab.yml): Full scanning/signing GHA workflow.", "• Speaker companions, lab manuals, and professional slide decks."]},
         
        # Slide 7: Timeline
        {"type": "content2", "title": "Development Timeline & Milestones",
         "c1_title": "Phases Complete", "c1_bullets": ["• Phase 1: Vulnerability Research (May 10 - May 14) - Exploit chain successfully compiled.", "• Phase 2: Runtime Remediation (May 15 - May 20) - Seccomp JSON and capability dropping filters verified."],
         "c2_title": "Integration & Delivery", "c2_bullets": ["• Phase 3: Pipeline Automation (May 21 - May 25) - GHA Trivy scanner and Cosign keyless signing fully integrated.", "• Phase 4: Project Delivery (May 26 - May 29) - Slide decks, Word manuals, and readout guides completed."]},
         
        # Slide 8: Methodology
        {"type": "content2", "title": "DevSecOps Research Methodology",
         "c1_title": "Threat Emulation Strategy", "c1_bullets": ["• Step 1: Model a vulnerable base container to analyze physical system borders.", "• Step 2: Compile a localized multi-threaded exploit to study kernel race condition behavior in page tables."],
         "c2_title": "Containment Design Pattern", "c2_bullets": ["• Step 3: Implement fine-grained OS-level controls (Seccomp syscall filters).", "• Step 4: Automate policy enforcement gates in CI/CD pipelines to cryptographically verify image authenticity."]},
         
        # Slide 9: Data/Results
        {"type": "content2", "title": "Technical Results & Containment Data",
         "c1_title": "Exploit Performance under Hardening", "c1_bullets": ["• Default Run: 100% success rate. Read-only targets overwritten in less than 5 seconds.", "• Seccomp Run: 0% success rate. The kernel blocks the 'madvise' call instantly, returning EPERM.", "• Exploit completely neutralized."],
         "c2_title": "Supply Chain Audit Data", "c2_bullets": ["• Trivy Scan Coverage: 100% of layers verified. CVE-2016-5195 is isolated.", "• Cosign Provenance: 100% of images verified. Short-lived 10-minute certificates logged securely in Rekor public ledger."]},
         
        # Slide 10: User/Stakeholder Impact
        {"type": "content2", "title": "Developer Experience & Operations Impact",
         "c1_title": "Zero-Friction DevSecOps Integration", "c1_bullets": ["• Before Hardening: Massive manual overhead. Fragile private key management in CI/CD, creating key leak vulnerabilities.", "• After Hardening: No manual private key storage. Fulcio CA dynamically signs code using automated GHA identity."],
         "c2_title": "System Stability Benefits", "c2_bullets": ["• Restricting capabilities and applying Seccomp profiles introduces zero performance overhead.", "• Security controls function purely as low-level kernel rules, securing the cluster seamlessly."]},
         
        # Slide 11: Risks & Mitigations
        {"type": "content2", "title": "Risk Register & Enterprise Mitigation",
         "c1_title": "Identified Risk Scenarios", "c1_bullets": ["• Risk 1: Host kernel suffers a new unpatched virtual memory zero-day exploit.", "• Risk 2: Public Fulcio/Rekor Sigstore services suffer an unexpected cloud outage.", "• Risk 3: Compromised developer credentials inside GitHub Action workspaces."],
         "c2_title": "SME Mitigations & Owners", "c2_bullets": ["• Mitigation 1: Deploy user namespace mapping and micro-VM container runtimes (Kata) to secure kernels. (SME)", "• Mitigation 2: Cache verified OIDC signatures locally in the cluster. (DevOps)", "• Mitigation 3: Enforce strict branch protections and environment secrets. (Admin)"]},
         
        # Slide 12: Budget/Resource Summary
        {"type": "content2", "title": "Budget, Resources & Cost Efficiency",
         "c1_title": "Zero-Cost Open-Source Stack", "c1_bullets": ["• Container Engine: Docker CE (Open-Source) $\rightarrow$ Planned: $0 | Actual: $0.", "• Vulnerability Scanner: Aquasec Trivy (Open-Source) $\rightarrow$ Planned: $0 | Actual: $0.", "• Cryptographic Signer: Cosign/Sigstore (Open-Source) $\rightarrow$ Planned: $0 | Actual: $0."],
         "c2_title": "Resource Allocation", "c2_bullets": ["• Build Runner: GitHub Actions Hosted Ubuntu Runner $\rightarrow$ Covered under free tier allocations.", "• High-performance security infrastructure built using 100% industry-standard open-source components."]},
         
        # Slide 13: Lessons Learned
        {"type": "content2", "title": "Actionable Lessons & Key Insights",
         "c1_title": "Container Boundaries are Logical", "c1_bullets": ["• Lesson 1: Namespaces and cgroups do not provide hardware physical boundaries. They are logical kernel parameters.", "• Lesson 2: Shared kernel models demand prompt host-level patching schedules."],
         "c2_title": "Signing and Scanning Synergy", "c2_bullets": ["• Lesson 3: Signing a vulnerable image is security theater. Static Trivy scanning must precede Cosign signing.", "• Lesson 4: Restricting syscalls (Seccomp) is highly effective at stopping exploits before they hit the kernel."]},
         
        # Slide 14: Next Steps & Asks
        {"type": "content2", "title": "Future Roadmap & Strategic Decisions",
         "c1_title": "Upcoming Security Phases", "c1_bullets": ["• Step 1: Transition application pods from standard containers to gVisor user-space kernel runtimes.", "• Step 2: Implement Kyverno Validating Admission Controllers inside the Kubernetes cluster."],
         "c2_title": "Decisions Needed & Owners", "c2_bullets": ["• Approve the deployment of user-space sandboxing in multi-tenant environments. (IITJ Staff / CISO)", "• Timeline: Target full admission controller enforcement by Q3 2026."]},
         
        # Slide 15: Appendix
        {"type": "content2", "title": "Appendix & Technical Reference Library",
         "c1_title": "Primary Code Assets", "c1_bullets": ["• Exploit: [dirtyc0w.c](file:///d:/Dirty%20cow%20Docker/dirtyc0w.c) - maps MAP_PRIVATE, calls madvise and writes /proc/self/mem.", "• Profile: [hexaforce-seccomp.json](file:///d:/Dirty%20cow%20Docker/hexaforce-seccomp.json) - blocks madvise system call."],
         "c2_title": "Automation & Guides", "c2_bullets": ["• Pipeline: [.github/workflows/dirtycow-lab.yml](file:///d:/Dirty%20cow%20Docker/.github/workflows/dirtycow-lab.yml).", "• Handbooks: [Hexa_Force_Secure_Lab_Guide.docx](file:///d:/Dirty%20cow%20Docker/Hexa_Force_Secure_Lab_Guide.docx).", "• Readout: [Hexa_Force_Presentation_Readout_Guide.docx](file:///d:/Dirty%20cow%20Docker/Hexa_Force_Presentation_Readout_Guide.docx)."]}
    ]
    
    # ----------------------------------------------------
    # Slide Speaker Notes Mapped Slide-by-Slide (Slides 2-14)
    # ----------------------------------------------------
    speaker_notes = [
        # Slide 1
        "",
        # Slide 2 (Exec Summary)
        "Good morning, members of the committee. Slide 2 presents our executive summary: we have successfully built and hardened a 4-stage container exploit chain. By dropping capabilities, running read-only filesystems, and implementing keyless Cosign signing via OIDC, we achieved 100% containment of the host-kernel privilege escalation threat.",
        # Slide 3 (Context)
        "On Slide 3, we analyze the core problem: standard containers share the host Linux kernel, creating a single point of failure. If the host kernel has a privilege escalation vulnerability like Dirty COW, the logical isolation of namespaces collapses, allowing low-privileged container processes to escape to the host.",
        # Slide 4 (KPIs)
        "Slide 4 displays our key performance indicators: we successfully reduced the host breakout escape likelihood from 100% vulnerable to 0% hardened. We also achieved 100% automated scanning with Trivy and 100% cryptographic signing using Cosign OIDC tokens.",
        # Slide 5 (Solution)
        "Moving to Slide 5, we present our DevSecOps solution architecture: it combines pipeline build gates (Trivy scans and Cosign keyless signing) with host runtime hardening, notably employing a custom Seccomp profile to block the vulnerable madvise system call.",
        # Slide 6 (Scope)
        "Slide 6 maps the laboratory deliverables: we have produced a hardened Dockerfile, the C exploit code, the custom Seccomp profile, and the automated GHA pipeline, along with comprehensive speaker notes and guides.",
        # Slide 7 (Timeline)
        "Slide 7 outlines our timeline: research and exploit modeling were finalized in mid-May, followed by runtime remediation and pipeline automation. We are now delivering all refined presentations and manuals.",
        # Slide 8 (Methodology)
        "On Slide 8, we outline our research methodology: we model vulnerable environments to isolate threat pathways, compile localized thread exploit races to analyze page-table behavior, and then deploy system-level and pipeline-level containment.",
        # Slide 9 (Results)
        "Slide 9 highlights our technical results: the default exploit overwrites read-only targets in under 5 seconds, but under our custom Seccomp profile, the exploit fails instantly. Furthermore, our Trivy scans successfully analyzed 100% of container layers.",
        # Slide 10 (DX Impact)
        "Slide 10 discusses the developer experience: by migrating to keyless Sigstore signing, we eliminated the fragile overhead of managing private keys in CI/CD environments, achieving automated cryptographic trust without security friction.",
        # Slide 11 (Risks)
        "On Slide 11, we evaluate system risks: key risks include host kernel zero-days and OIDC provider outages. We mitigate these by caching verified signatures locally and planning transitions to user-space kernels.",
        # Slide 12 (Budget)
        "Slide 12 summarizes the budget: by leveraging industry-standard open-source components such as Docker, Trivy, and Cosign, we built a highly robust enterprise security suite at zero cost.",
        # Slide 13 (Lessons Learned)
        "Slide 13 lists our critical lessons: container namespace boundaries are purely logical parameters, and shared host kernels demand prompt patching. Synergizing image scanning with syscall blocking provides robust defense-in-depth.",
        # Slide 14 (Next Steps)
        "Finally, Slide 14 presents our future roadmap: our next steps are to evaluate gVisor sandboxing in multi-tenant environments and deploy Kyverno Admission Controllers to enforce container signatures in our Kubernetes cluster."
    ]
    
    # Generate Slides
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        apply_dark_bg(slide)
        
        if slide_data["type"] == "title":
            add_title_slide(slide, slide_data["title"], slide_data["subtitle"], slide_data["cat"])
        else:
            add_header(slide, slide_data["title"])
            
            # Draw Two Columns
            col_w = Inches(5.6)
            col_h = Inches(4.8)
            top_y = Inches(1.8)
            
            # Left Card
            draw_card(slide, Inches(0.75), top_y, col_w, col_h)
            add_card_text(slide, Inches(0.75), top_y, col_w, col_h, slide_data["c1_title"], slide_data["c1_bullets"], title_color=C_PRIMARY)
            
            # Right Card (If hardening, outline with Emerald Green)
            border_c = C_GREEN if "Hardening" in slide_data["title"] or "Remediation" in slide_data["title"] or "Results" in slide_data["title"] or "Takeaways" in slide_data["title"] or "Outcomes" in slide_data["title"] else None
            draw_card(slide, Inches(6.8), top_y, col_w, col_h, border_color=border_c)
            add_card_text(slide, Inches(6.8), top_y, col_w, col_h, slide_data["c2_title"], slide_data["c2_bullets"], title_color=C_GREEN if border_c else C_PRIMARY)
            
        # Inject Speaker Notes
        if i < len(speaker_notes):
            notes_slide = slide.notes_slide
            tf_notes = notes_slide.notes_text_frame
            tf_notes.text = speaker_notes[i]
            
    filename = "FINAL briefing.pptx"
    prs.save(filename)
    print(f"SUCCESS: Rebuilt {filename} successfully with 15 slides and embedded speaker notes!")

if __name__ == "__main__":
    build_final_briefing()
