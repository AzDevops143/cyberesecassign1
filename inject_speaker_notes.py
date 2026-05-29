import sys
from pptx import Presentation

# ----------------------------------------------------
# DECK 1: Hexa Force Secure Lab Presentation Speaker Notes
# ----------------------------------------------------
deck1_notes = [
    # Slide 1
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Good morning, members of the evaluation committee and distinguished guests. Welcome to the technical delivery "
        "of the Hexa Force Container Security Lab. Today, we are presenting our end-to-end container security demonstration, "
        "modeling a high-fidelity 4-stage attack chain and corresponding defense-in-depth hardening strategies. Over the "
        "next few minutes, we will walk you through how a local privilege escalation vulnerability can cascade into complete "
        "host takeover, and how modern cryptographic verification using keyless Sigstore/Cosign signing with GitHub Actions "
        "OIDC robustly secures the software delivery pipeline. Let us begin by reviewing our core objectives.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Establish the Hexa Force brand immediately with an authoritative and professional delivery.\n"
        "• Set the scope of the presentation: it is a balanced demonstration of an advanced attack chain and enterprise-grade cryptographic defense.\n"
        "• Emphasize key modern pipeline terms: Sigstore, Cosign, OIDC federation, and GHCR registry integrations.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What motivated the choice of this specific attack chain?\n"
        "Answer: We chose this specific progression because it realistically models how simple runtime misconfigurations and unpatched host-level vulnerabilities interact in enterprise microservices. It demonstrates that container security cannot rely on runtime isolation alone, but must encompass host hardening, image provenance, and CI/CD integrity."
    ),
    # Slide 2
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Our presentation today is structured into four main areas. First, we establish the threat landscape and the "
        "4 Security Pillars of Containers. Next, we walk step-by-step through our 4-stage attack simulation, showcasing how "
        "we break kernel isolation, processes, APIs, and volume mounts. Third, we present our cryptographic supply chain "
        "integrity pipeline using Cosign and OpenID Connect. Finally, we conclude with our enterprise-grade hardening "
        "recommendations and DevSecOps best practices to secure any modern production cluster.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Guide the audience through the agenda to show structural clarity and excellent preparation.\n"
        "• Emphasize that the lab moves logically from threat identification (simulation) to threat neutralization (hardening).\n"
        "• Make it clear that the defense-in-depth framework applies to all levels: runtime, pipeline, and host kernel.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the most critical phase of the agenda?\n"
        "Answer: While the exploit phase demonstrates the weakness, the supply chain integrity section is the most critical. Securing the container runtime is only half the battle; we must also guarantee code provenance through modern cryptographic signing to prevent unauthorized container execution."
    ),
    # Slide 3
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"To secure any containerized environment, we must evaluate the four foundational security pillars: (1) Host Kernel "
        "Isolation, (2) Process & Capabilities Isolation, (3) Daemon & API Security, and (4) Volume & Filesystem Security. "
        "Standard container systems share the host kernel, which means a kernel vulnerability exposes the host. Processes "
        "inside a container must be strictly limited via Linux capabilities to prevent privilege abuse. The Docker API daemon "
        "socket must never be exposed to untrusted environments, and physical volume mounts must be strictly governed to "
        "prevent directory traversal and host escapes. Let us look at Stage 1, where kernel isolation is broken.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Introduce the 4 Pillars as the industry-standard frame of reference for evaluating container security.\n"
        "• Explain the shared kernel model: containers are lightweight because they share the host kernel, which is also their primary single point of failure.\n"
        "• Emphasize that security must be applied across all 4 pillars simultaneously; securing only one layer leaves the system vulnerable.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why is kernel sharing the ultimate double-edged sword for containers?\n"
        "Answer: Kernel sharing provides the lightweight, high-performance execution that makes containers so successful compared to heavy virtual machines. However, it also introduces a single point of failure: any local privilege escalation vulnerability in the host's Linux kernel (like Dirty COW) is immediately shared by every container on that host, collapsing the primary isolation boundary."
    ),
    # Slide 4
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"At Stage 1, we execute our first exploit inside an unprivileged container. Running as the unprivileged user "
        "'victim', we execute our compiled binary, 'dirtyc0w'. This exploit targets the Dirty COW kernel vulnerability, "
        "which is a race condition in the kernel's virtual memory subsystem's handling of Copy-on-Write. We target a "
        "read-only secret file, '/var/secret/target.txt'. By using madvise with the MADV_DONTNEED flag, we race the kernel's "
        "write mechanisms, tricking the kernel into writing root-owned data into a read-only memory segment. This instantly "
        "elevates our privilege from an unprivileged user to container root.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Focus heavily on the virtual memory mechanics: Copy-on-Write, madvise system call, and the O_RDONLY file bypass.\n"
        "• Clearly declare CVE-2016-5195 to demonstrate deep technical mastery and precise vulnerability tracking.\n"
        "• Emphasize that this first exploit is executed as a non-root, unprivileged user, simulating a compromised microservice shell.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why does a local privilege escalation inside a container matter if it is isolated?\n"
        "Answer: Local privilege escalation is the beachhead. While namespaces still limit what a container root user can see, obtaining root privileges is the prerequisite for abusing other exposed interfaces, such as Linux capabilities (like CAP_SYS_ADMIN), raw network sockets, or system files."
    ),
    # Slide 5
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"In Stage 2, we exploit process-level privileges. Having achieved container root in Stage 1, we now check "
        "the container's capabilities. If a container is run with elevated Linux capabilities, such as CAP_SYS_ADMIN, "
        "the boundaries of standard container isolation begin to disintegrate. Linux capabilities partition traditional "
        "superuser root powers into distinct, assignable units. By abusing CAP_SYS_ADMIN, we gain the authority to modify "
        "kernel namespaces and mount external devices, laying the physical groundwork for our escape to the host filesystem.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Explain the concept of Linux capabilities: splitting traditional all-or-nothing root powers into fine-grained permissions.\n"
        "• Emphasize CAP_SYS_ADMIN as the root of all process isolation failures inside containers.\n"
        "• Discuss the principle of least privilege: containers must drop all unnecessary capabilities by default.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What are some alternatives to granting CAP_SYS_ADMIN to containers?\n"
        "Answer: Administrators should perform strict capability profiling and grant only the narrowest possible capabilities (like CAP_NET_BIND_SERVICE for low-port binding) instead of CAP_SYS_ADMIN. If system-level tasks are absolutely required, a dedicated daemon or sidecar pattern should be evaluated."
    ),
    # Slide 6
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Stage 3 demonstrates a critical lateral movement vector: the exposed Docker socket. In many development "
        "environments, administrators mount the host's Docker socket—'/var/run/docker.sock'—directly into the container "
        "to enable nested docker tasks. Since our Stage 1 exploit gave us root privileges inside the container, we have "
        "full read/write permission to this socket. Because the Docker socket is the direct command API to the host's "
        "Docker engine, we can now send API requests directly to the host daemon, giving us full command over the host's "
        "container lifecycle.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Highlight the absolute danger of the Docker socket (/var/run/docker.sock) exposure.\n"
        "• Describe the Docker socket as the entry point to the host's physical container engine.\n"
        "• Explain how controlling the socket allows the containerized root user to send administrative API commands directly to the host daemon.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the principal security risk of mounting the Docker socket inside a container?\n"
        "Answer: Mounting the Docker socket is equivalent to granting root access to the physical host. Anyone with access to the socket can send commands to the host Docker daemon to pull arbitrary images, create new privileged containers, and mount the host's physical hard drives, rendering container isolation completely obsolete."
    ),
    # Slide 7
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"In Stage 4, we execute the final breakout. Using our access to the Docker socket, we issue an API command to "
        "launch a new, highly privileged container on the host. In its configuration, we mount the host's physical root "
        "filesystem to '/mnt/host' inside the container. Because this container runs with full host namespace access, "
        "we traverse directly into '/mnt/host/etc/cron.d' and write a malicious cron job. When the host's cron daemon "
        "executes, it runs our script and opens a root-privileged reverse shell back to our attack machine. We have achieved "
        "complete host takeover.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Detail the mechanical steps of host escape: privileged container instantiation $\rightarrow$ root filesystem mount $\rightarrow$ cron job injection.\n"
        "• Emphasize the threat of persistence: writing to `/etc/cron.d` guarantees that the host daemon will re-execute the reverse shell periodically.\n"
        "• Establish that this completes a full, high-severity exploit chain from a low-privileged container user to total physical host control.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How can host filesystem mounting be blocked at the enterprise level?\n"
        "Answer: Enterprises should implement Kubernetes Pod Security Admissions (PSA) or open-source policies (like OPA Gatekeeper or Kyverno) to strictly enforce that containers cannot run in privileged mode, share the host PID/IPC namespaces, or mount host paths."
    ),
    # Slide 8
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Beyond runtime hardening, Hexa Force secures the software supply chain using state-of-the-art keyless "
        "container signing via Cosign. Traditional signing relies on storing highly sensitive cryptographic private keys. "
        "If these keys are leaked, the supply chain is compromised. Hexa Force avoids this completely. We leverage OpenID "
        "Connect to exchange GitHub Action build identities for temporary, 10-minute certificates. Cosign signs the built "
        "image, recording the signature in the immutable public Rekor transparency log, guaranteeing container origin and integrity.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Contrast traditional key management with modern keyless container signing.\n"
        "• Define the role of Sigstore: Fulcio issues short-lived certificates, Rekor provides the immutable public ledger.\n"
        "• Explain the cryptographic trust path: the signature is bound to a specific OIDC builder identity, verifying code provenance.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is Fulcio's role in keyless signing?\n"
        "Answer: Fulcio is the Sigstore Certificate Authority. It issues short-lived (10-minute) X.509 certificates. The certificate is bound to an OIDC token issued by a trusted identity provider (like GitHub Actions), proving that the build pipeline itself signed the container."
    ),
    # Slide 9
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"We operationalize this supply chain defense by automating the entire pipeline inside GitHub Actions. Upon a push "
        "to the main branch, our workflow triggers. It builds the Docker image, dynamically converts the image tag to "
        "lowercase using bash expansion to ensure compatibility with container registry specifications, and pushes it to "
        "the GitHub Container Registry. Immediately, Cosign uses the OIDC identity to sign the image. This seamless, "
        "automated pipeline guarantees that only verified and signed code is allowed to run in our production environment.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Highlight the automation aspects of the GHA workflow: building, lowercasing, pushing to GHCR, and signing.\n"
        "• Explain the lowercase string translation `${GITHUB_REPOSITORY,,}` as a crucial Docker OCI registry compliance step.\n"
        "• Position the CI/CD pipeline as the core policy enforcement point of your security program.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How does signature verification work inside a Kubernetes cluster?\n"
        "Answer: A validation controller (like Cosign Policy Controller or Kyverno) runs as an Admission Controller in the cluster. When a new pod request is made, the controller queries the Rekor transparency log to verify the image's cryptographic signature against the allowed OIDC builder identity, rejecting any unsigned or unverified images."
    ),
    # Slide 10
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Let us examine how we completely remediate this threat. Hexa Force implements a comprehensive Defense-in-Depth "
        "architecture. First, we restrict container capabilities by explicitly dropping CAP_SYS_ADMIN. Second, we run the "
        "container as a non-root user and mark the root filesystem as read-only. Third, we block all exposure of the Docker "
        "socket. Finally, we apply kernel security patches to eliminate Dirty COW at the host level. Even if an attacker "
        "gains an initial shell, these combined controls prevent them from ever gaining root or escaping the container.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Emphasize the four layers of defense: capability dropping, non-root runtimes, read-only rootfs, and kernel patching.\n"
        "• Discuss the concept of Defense-in-Depth: even if one layer fails, subsequent layers prevent the exploit from succeeding.\n"
        "• Position these controls as standard, low-overhead enterprise recommendations.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the security value of a read-only container root filesystem?\n"
        "Answer: A read-only root filesystem prevents an attacker from writing files, installing tools, downloading exploit scripts, or compiling C binaries (like `dirtyc0w.c`) inside the container. It dramatically reduces the available attack surface and prevents runtime environmental tampering."
    ),
    # Slide 11
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"In conclusion, the Hexa Force Container Security Lab successfully demonstrates the extreme vulnerability of "
        "default container runtimes, while delivering an enterprise-ready blueprint for comprehensive defense. By pairing "
        "strict runtime hardening policies with cryptographic supply chain verification, we ensure complete resilience "
        "against host compromises and supply chain injection. Thank you for your time. The Hexa Force team is now open "
        "to your technical questions.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Deliver a clean, professional closing statement summarizing the core achievements.\n"
        "• Reiterate the value of combining runtime isolation with cryptographic verification.\n"
        "• Invite questions with absolute confidence, ready to leverage your detailed Q&As.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the single most important lesson from this research?\n"
        "Answer: Security is a multi-layered process. We cannot rely on container virtualization as a bulletproof sandbox. We must harden the host kernel, enforce the principle of least privilege inside the container, and use cryptographic pipelines to guarantee that only trusted code ever enters the execution environment."
    )
]

# ----------------------------------------------------
# DECK 2: Dirty COW Technical Briefing Speaker Notes
# ----------------------------------------------------
deck2_notes = [
    # Slide 1
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Welcome to Part 2 of our briefing: The Dirty COW Technical Deep Dive. In this session, we will demystify the "
        "precise mechanisms behind CVE-2016-5195, commonly known as Dirty COW. This kernel exploit bypassed decades of "
        "security policies in Linux systems by weaponizing a race condition inside the virtual memory management subsystem. "
        "We will explore how memory pages are mapped, how Copy-on-Write page cache mechanisms work, and how we can force "
        "the kernel to write user data directly to read-only root files. Let us dive into the exploit's background.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Introduce the briefing as a deep, academically rigorous analysis of a legendary kernel vulnerability.\n"
        "• Frame the presentation around core operating system structures: virtual memory, page fault handlers, and thread scheduling.\n"
        "• Clearly declare the scope: it covers the CVE profile, exploitation theory, code walkthrough, and kernel patch analysis.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why is it called \"Dirty COW\"?\n"
        "Answer: \"COW\" refers to the \"Copy-on-Write\" mechanism. \"Dirty\" refers to marking a page in physical memory as modified in the page tables, indicating that its changes must be written back to the disk."
    ),
    # Slide 2
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"The briefing today is structured to provide both foundational theory and a hands-on lab demonstration. "
        "We will start with memory management basics, move through the COW design and vulnerability analysis, "
        "perform a deep code review of dirtyc0w.c, showcase our hands-on Docker laboratory environment, "
        "analyze the real-world impact, and outline concrete remediation strategies.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Lay out the technical roadmap for the deep dive to establish structured, expert-level delivery.\n"
        "• Note that the briefing bridges theory (operating system architecture) with practice (compiled exploits in containers).\n"
        "• Emphasize that understanding the vulnerability code is key to understanding modern mitigation strategies.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why is a deep dive on a 2016 vulnerability still highly valuable today?\n"
        "Answer: Dirty COW is the canonical textbook example of a virtual memory race condition exploit. By understanding it, security engineers gain deep insights into how kernel memory mapping works and why container environments (which share the host kernel) are highly vulnerable to shared resource exploits."
    ),
    # Slide 3
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"To contextualize this threat, let's examine the profile of CVE-2016-5195. Discovered in 2016, this "
        "vulnerability carries a High CVSS score. It affected almost every Linux kernel version prior to October 2016. "
        "The beauty—and danger—of Dirty COW is that it requires absolutely no special system privileges to execute. Any "
        "local user, even an unprivileged process running inside a Docker container, can trigger this exploit to bypass "
        "file permissions and write root-owned configurations. Let us look at the memory architecture that made this possible.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Cite precise vulnerability metadata: CVE-2016-5195, High severity CVSS rating, and pre-2016 kernel versions.\n"
        "• Highlight that this is a local privilege escalation (LPE) vector, requiring zero administrative starting permissions.\n"
        "• Set up the core threat: it completely compromises local system file integrity.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Can Dirty COW be executed over the network without local access?\n"
        "Answer: No, it is strictly a Local Privilege Escalation (LPE) vulnerability. The attacker must already have local shell access or code execution capabilities on the system to compile and run the exploit binary."
    ),
    # Slide 4
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Let's review Linux virtual memory management. The operating system uses virtual memory addresses to isolate "
        "processes from each other and from physical hardware. Physical RAM is divided into fixed-size units called page "
        "frames, while virtual memory is managed in pages. The CPU's Memory Management Unit (MMU) and page tables translate "
        "virtual addresses to physical RAM addresses. Processes map files into their virtual memory using mmap, which allows "
        "reading or writing file data directly through RAM operations.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Explain the role of the MMU (Memory Management Unit) in performing address translations.\n"
        "• Define key terms: page frames (physical RAM), pages (virtual addresses), and page table entries.\n"
        "• Establish how `mmap` optimizes file reading by loading disk data directly into memory maps.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What role does the MMU play in memory protection?\n"
        "Answer: The MMU translates addresses and enforces page-level access permissions (e.g. read-only, read-write, executable) defined in the page table entries, throwing a segmentation fault if an unauthorized write is attempted."
    ),
    # Slide 5
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Copy-on-Write is a crucial memory optimization. When a process maps a file privately using the MAP_PRIVATE flag, "
        "the kernel does not create a duplicate copy of the file in RAM. Instead, it maps the process's virtual memory "
        "to the same physical page cache frame as the underlying file to save memory. When the process attempts to write "
        "to this page, the kernel detects the write, allocates a new physical page frame, copies the original page content, "
        "updates the process's page table to point to this new private copy, and executes the write. This is the Copy-on-Write process.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Explain Copy-on-Write (COW) as an elegant RAM optimization technique used universally in modern operating systems.\n"
        "• Define `MAP_PRIVATE`: it guarantees that write modifications are made to a private memory space, never touching the underlying file.\n"
        "• Describe the transition from shared physical pages to private copy allocation during a write operation.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How does MAP_PRIVATE protect the original backing file?\n"
        "Answer: MAP_PRIVATE guarantees that all write operations are executed against a private local copy allocated in physical RAM, ensuring that the original backing file on disk is never modified."
    ),
    # Slide 6
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"The core vulnerability is a non-atomic operation in the page fault handler. When a process writes to a private "
        "read-only page, the kernel follows a multi-step routine: (1) Find the physical page, (2) Verify write permissions "
        "(which fails, triggering a page fault), (3) Perform the Copy-on-Write allocation, and (4) Update page table entries. "
        "Because this process is not atomic, the kernel must yield or check permissions at multiple boundaries. If page "
        "table entries are cleared concurrently during these checks, the kernel gets confused, allowing a concurrent write "
        "to write directly to the original read-only page before the COW is completed.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Define \"non-atomic\": a sequence of operations that can be interrupted or preempted by other CPU instructions.\n"
        "• Describe the precise moment the vulnerability occurs: the gap between checking the write permission and updating the page table entry.\n"
        "• Explain how concurrent threads can corrupt this gap, causing the kernel to fall back to writing to the read-only backing page.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What does non-atomic mean in this context?\n"
        "Answer: It means the operation is broken into multiple individual CPU instructions. In a multi-threaded system, other threads can execute instructions in between these steps, altering the memory state before the operation is completed."
    ),
    # Slide 7
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"The exploit utilizes two key system calls: (1) mmap with the MAP_PRIVATE flag to map the target read-only "
        "file into virtual memory, and (2) madvise with the MADV_DONTNEED flag. The madvise system call tells the kernel "
        "that the process no longer needs the specified memory range. The kernel responds by discarding the physical page "
        "mappings, forcing the next read to load the data directly from the original backing file. Additionally, the "
        "exploit writes directly to `/proc/self/mem` to bypass standard MMU write restrictions.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Detail the roles of `mmap(MAP_PRIVATE)` and `madvise(MADV_DONTNEED)`.\n"
        "• Explain that `MADV_DONTNEED` forces the kernel to instantly invalidate the physical page tables mapping.\n"
        "• Introduce `/proc/self/mem` as the direct virtual process memory interface used to trigger the kernel-level write routing.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why is writing to `/proc/self/mem` used instead of direct pointer dereferencing?\n"
        "Answer: Direct pointer writes are blocked by the CPU's MMU immediately, raising a segmentation fault. Writing through the `/proc/self/mem` interface routes the write through the kernel's virtual memory system, forcing the kernel itself to handle the write and execute the vulnerable page fault routine."
    ),
    # Slide 8
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"The exploit sequence runs two concurrent threads. Thread A runs in a tight loop calling `madvise(MADV_DONTNEED)` "
        "on our mapped memory address, repeatedly throwing away any newly allocated COW pages. Thread B runs in a parallel "
        "loop writing our payload to `/proc/self/mem` at the same virtual address. As they execute millions of times per "
        "second across CPU cores, a race condition occurs. The kernel completes the write permission check but gets "
        "preempted by madvise discarding the page. The write then executes directly against the original, read-only physical "
        "backing page, overwriting the file on disk.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Detail the concurrent execution flow: Thread A (`madvise`) and Thread B (`/proc/self/mem` write) racing at high speed.\n"
        "• Explain the physical collision: the write is directed to the read-only backing cache instead of the private copy.\n"
        "• Emphasize that the exploit requires multiple iterations in a loop to guarantee hitting the microsecond race window.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How many CPU cores are required to exploit the race condition?\n"
        "Answer: While it is theoretically possible on a single core via context switching, it is highly efficient on multi-core systems where Thread A and Thread B can run simultaneously on separate physical cores, maximizing the probability of racing the kernel's instructions."
    ),
    # Slide 9
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Dirty COW is particularly catastrophic in Docker container environments. Because Docker containers share the "
        "host's Linux kernel, any container—even one running as a low-privileged user—can run `dirtyc0w` to exploit the "
        "shared kernel. The exploit can bypass container isolation barriers entirely. An attacker inside a container "
        "can overwrite a read-only system file on the host's kernel or overwrite physical files shared between containers, "
        "collapsing the foundational promise of container isolation.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Explain why shared kernels represent the central security vulnerability of containerization.\n"
        "• Discuss the collapse of container isolation boundaries when a container process can overwrite host-level system files.\n"
        "• Reiterate that namespaces and cgroups do not prevent physical page cache corruption because both container and host share the same kernel cache.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How does Dirty COW bypass Docker's read-only file mappings?\n"
        "Answer: Docker's read-only mappings are enforced at the mount and namespace levels. However, because the exploit operates at the shared kernel virtual memory page level, it writes directly to the shared physical page cache. The kernel executes the physical write, rendering namespace and mount restrictions completely ineffective."
    ),
    # Slide 10
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Let's trace the threat pathway from privilege escalation to complete container escape. First, the attacker "
        "exploits an application vulnerability to gain a shell as a low-privileged container user. Second, they run "
        "Dirty COW to overwrite a root-owned file (like `/etc/passwd` or `/var/secret/target.txt`), gaining root "
        "privileges inside the container. Third, they leverage container root to access administrative capabilities "
        "or mounts (like the exposed `/var/run/docker.sock`), allowing them to launch new privileged containers, mount "
        "the physical host filesystem, and achieve complete host escape.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Illustrate the escalate-and-escape pathway step-by-step.\n"
        "• Connect Dirty COW (privilege escalation) with local boundary breakages (capabilities and exposed sockets).\n"
        "• Emphasize that a container compromise almost always leads to a host compromise if runtime capabilities are left unchecked.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Is container escape possible without local privilege escalation?\n"
        "Answer: Yes, if the container is already run with high privileges (like `--privileged` or mounting the host filesystem directly). However, if the container is properly locked down, privilege escalation (like Dirty COW) is the necessary first step to unlock the administrative APIs required for escape."
    ),
    # Slide 11
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Our C implementation `dirtyc0w.c` is elegant and compact. We open the target file, retrieve its size, and map "
        "it read-only but privately using `mmap`. We then spin up two parallel threads using `pthread_create`. Thread 1 "
        "runs `madviseThread` calling `madvise(map, 100, MADV_DONTNEED)` in a loop of 100,000 iterations. Thread 2 runs "
        "`procSelfMemThread` opening `/proc/self/mem`, seeking to the mapped address, and writing the payload in a loop of "
        "100,000 iterations. The parallel execution guarantees that the race condition will be triggered within seconds.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Walk through the core functions of `dirtyc0w.c`: `mmap`, `pthread_create`, `madviseThread`, and `procSelfMemThread`.\n"
        "• Explain the seek operation `lseek(f, (uintptr_t)map, SEEK_SET)` as the mapping pointer address offset aligner.\n"
        "• Detail the iteration scale (100,000 runs) that ensures the thread collision occurs quickly.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What happens to the target file if its size is smaller than the payload?\n"
        "Answer: The exploit can only overwrite data up to the original size of the mapped file. Attempting to write beyond the file boundary will fail or result in a bus error, which is why the exploit is typically used to overwrite configuration files (like `/etc/passwd`) where the text can be swapped in-place."
    ),
    # Slide 12
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"To demonstrate this exploit in a controlled environment, Hexa Force built a complete Docker-based laboratory. "
        "The laboratory uses a custom Dockerfile that copies `dirtyc0w.c` and a shell demo script `run_demo.sh` into an "
        "unprivileged container running under the user `victim`. A read-only file `/var/secret/target.txt` is created "
        "inside the container as the target. The environment is hosted on an unpatched Linux kernel to allow students "
        "and security engineers to safely observe the exploit in real-time.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Describe the laboratory sandbox components: unprivileged container environment, `victim` user, and read-only target file.\n"
        "• Explain that the container is run on an unpatched Linux kernel within a nested hypervisor or isolated virtual machine.\n"
        "• Emphasize that the lab design replicates real-world microservice exploitation in a safe, controlled manner.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How do you prevent students from accidentally compromising the host in this lab?\n"
        "Answer: The lab is run in an isolated virtual machine or sandbox environment. Although the container can theoretically exploit the kernel to escape, running the entire demonstration inside a dedicated test VM guarantees that any host compromise is contained within the sandbox."
    ),
    # Slide 13
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Let's examine the outcome of executing the lab. When we run `./run_demo.sh`, the script compiles `dirtyc0w.c` "
        "using gcc and executes the binary against the read-only file `/var/secret/target.txt` with a replacement payload. "
        "The threads spin up, complete their 100,000 iterations, and print their completion. When we read `/var/secret/target.txt` "
        "immediately after, we find that the read-only file has been successfully overwritten, proving the exploit's success in real-time.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Outline the real-time execution steps: gcc compilation, target mapping, multi-threaded racing, and outcome reading.\n"
        "• Emphasize the immediate, visual feedback of the overwrite to prove exploit success.\n"
        "• Position the lab as a concrete proof of the kernel's virtual memory vulnerability under high-speed thread competition.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Why does the exploit sometimes require multiple executions to succeed?\n"
        "Answer: Because it relies on a timing race condition. Depending on CPU scheduling, system load, and thread execution alignment, the race window may not be hit immediately. However, running the loops 100,000 times almost always guarantees success within a single run."
    ),
    # Slide 14
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"The real-world impact of Dirty COW was massive. It affected millions of Linux servers, Android devices, "
        "and cloud microservice environments worldwide. In multi-tenant environments—where multiple customers share the "
        "same physical server kernel—Dirty COW allowed malicious tenants to bypass isolation and access other tenants' data. "
        "It serves as a stark reminder of why prompt kernel patching is a fundamental requirement of cloud infrastructure security.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Discuss the global scope of CVE-2016-5195: millions of servers, Android mobile devices, and virtualization nodes compromised.\n"
        "• Explain the risk to multi-tenant cloud hosting: unprivileged users reading or writing adjacent customer memory segments.\n"
        "• Reiterate that kernel level patching is not an optional operation, but a foundational cloud security mandate.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: Was Dirty COW ever exploited in the wild?\n"
        "Answer: Yes. Security researchers and threat intelligence firms detected in-the-wild exploits using Dirty COW to root Android devices, bypass mobile security controls, and compromise multi-tenant web hosting servers before patches were universally applied."
    ),
    # Slide 15
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"Remediating Dirty COW requires a multi-layered security strategy. The primary defense is applying kernel updates "
        "to patch the vulnerability. Beyond kernel patching, we must implement Defense-in-Depth. In containers, we should "
        "strictly drop administrative capabilities (like CAP_SYS_ADMIN), run processes as non-root users, utilize read-only "
        "root filesystems, and employ AppArmor or SELinux profiles to restrict system calls like madvise or write access "
        "to `/proc/self/mem`.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Position the core solution: apply the Linux kernel security patch that fixes Copy-on-Write page table routines.\n"
        "• Walk through the supplemental Defense-in-Depth layers: capability dropping, non-root runtimes, and read-only filesystems.\n"
        "• Explain how LSMs (Linux Security Modules) like AppArmor or SELinux can block unauthorized syscalls like `madvise`.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: How can AppArmor or SELinux mitigate Dirty COW if the kernel is unpatched?\n"
        "Answer: AppArmor or SELinux can block the process from writing to `/proc/self/mem` or restrict the execution of specific system calls, neutralising the exploit's primary vector even if the underlying kernel vulnerability is still present."
    ),
    # Slide 16
    (
        "=== SPOKEN TALK TRACK ===\n"
        "\"In conclusion, the Dirty COW technical briefing highlights the deep architectural risks of shared-kernel "
        "virtualization. While containerization offers unparalleled efficiency, it demands rigorous host-level patching, "
        "strict runtime isolation rules, and continuous security automation. The Hexa Force DevSecOps architecture "
        "successfully mitigates these risks, delivering a robust blueprint for modern enterprise security. Thank you, "
        "and we are open to any questions.\"\n\n"
        "=== CRITICAL TALKING POINTS ===\n"
        "• Deliver a strong technical summary linking shared-kernel realities to continuous security patching.\n"
        "• Reiterate that modern microservices demand runtime sandboxing, supply chain integrity, and regular kernel audits.\n"
        "• Open the floor to final technical questions with absolute poise and authority.\n\n"
        "=== ANTICIPATED EVALUATOR Q&A ===\n"
        "Question: What is the most important lesson from the Dirty COW briefing?\n"
        "Answer: We must never treat containers as full security boundaries by default. Container security requires a continuous DevSecOps pipeline combining host kernel patching, least-privilege runtimes, and supply chain integrity validation to guarantee absolute system security."
    )
]

def inject_notes(filename, notes_data):
    try:
        prs = Presentation(filename)
        if len(prs.slides) != len(notes_data):
            print(f"ERROR: Slide count mismatch for {filename}. Slide count in file: {len(prs.slides)}, data length: {len(notes_data)}")
            return False
        for i, slide in enumerate(prs.slides):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes_data[i]
        prs.save(filename)
        print(f"SUCCESS: Injected speaker notes into {filename}")
        return True
    except Exception as e:
        print(f"ERROR: Failed to inject notes into {filename}. Reason: {e}")
        return False

def main():
    print("[*] Starting slide notes injection...")
    
    # Inject Deck 1
    d1_file = "Hexa_Force_Secure_Lab_Presentation.pptx"
    s1 = inject_notes(d1_file, deck1_notes)
    
    # Inject Deck 2
    d2_file = "Dirty_COW_Technical_Briefing.pptx"
    s2 = inject_notes(d2_file, deck2_notes)
    
    if s1 and s2:
        print("[+] SUCCESS: All speaker notes successfully embedded directly inside slide files!")
    else:
        print("[-] WARNING: One or more injections encountered issues.")

if __name__ == "__main__":
    main()
